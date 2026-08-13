from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from PIL import Image

ROOT = Path(__file__).resolve().parents[1]
GAME = ROOT / "Game"
OUT = ROOT / "presentation" / "Local_Arcade_Portfolio.pptx"

# Editorial portfolio palette. Most slides use paper-like backgrounds and
# restrained blue accents so screenshots and engineering explanations lead.
BG = RGBColor(246, 247, 249)
PANEL = RGBColor(255, 255, 255)
PANEL2 = RGBColor(238, 241, 245)
WHITE = RGBColor(31, 41, 55)
MUTED = RGBColor(91, 103, 120)
MINT = RGBColor(39, 111, 151)
PINK = RGBColor(73, 103, 145)
YELLOW = RGBColor(142, 105, 54)
BLUE = RGBColor(49, 94, 154)
LINE = RGBColor(210, 216, 224)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank = prs.slide_layouts[6]

def rect(slide, x, y, w, h, fill=PANEL, line=None, radius=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                                   Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line or fill
    return shape

def text(slide, value, x, y, w, h, size=18, color=WHITE, bold=False,
         font="Malgun Gothic", align=PP_ALIGN.LEFT, valign=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True; tf.vertical_anchor = valign
    p = tf.paragraphs[0]; p.text = value; p.alignment = align
    p.font.name = font; p.font.size = Pt(size); p.font.bold = bold; p.font.color.rgb = color
    return box

def rich_lines(slide, lines, x, y, w, h, size=16, gap=8):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame; tf.clear(); tf.word_wrap = True
    for i, (lead, body, color) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap); p.font.name = "Malgun Gothic"; p.font.size = Pt(size); p.font.color.rgb = WHITE
        r = p.add_run(); r.text = lead; r.font.bold = True; r.font.color.rgb = color
        r = p.add_run(); r.text = body; r.font.color.rgb = WHITE
    return box

def base(title, kicker, number):
    s = prs.slides.add_slide(blank); s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
    text(s, kicker, .62, .30, 6, .25, 9, MUTED, False, "Arial")
    text(s, title, .62, .62, 11.9, .55, 25, WHITE, True)
    rect(s, .55, 1.19, 12.2, .015, LINE, LINE, False)
    text(s, f"Local Arcade  |  {number:02d}", 10.65, 7.12, 2.1, .2, 8, MUTED, False, "Arial", PP_ALIGN.RIGHT)
    return s

def add_image(slide, path, x, y, w, h, crop_top=0, mode="contain"):
    path = Path(path)
    with Image.open(path) as im:
        iw, ih = im.size
    effective_h = ih * (1-crop_top) if crop_top else ih
    target_ratio = w/h; image_ratio = iw/effective_h
    if mode == "contain":
        # Preserve the full screenshot. Empty space is intentional and makes UI
        # context easier to understand than aggressive cropping.
        rect(slide, x, y, w, h, PANEL2, LINE, False)
        if image_ratio > target_ratio:
            ph = w / image_ratio
            px, py, pw = x, y + (h-ph)/2, w
        else:
            pw = h * image_ratio
            px, py, ph = x + (w-pw)/2, y, h
        pic = slide.shapes.add_picture(str(path), Inches(px), Inches(py), width=Inches(pw), height=Inches(ph))
    else:
        pic = slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(w), height=Inches(h))
        if crop_top: pic.crop_top = crop_top
        if image_ratio > target_ratio:
            extra = 1 - target_ratio/image_ratio; pic.crop_left = extra/2; pic.crop_right = extra/2
        elif image_ratio < target_ratio:
            extra = 1 - image_ratio/target_ratio; pic.crop_top = max(pic.crop_top, extra/2); pic.crop_bottom = extra/2
    pic.line.color.rgb = LINE
    return pic

def code(slide, value, x, y, w, h, title="IMPLEMENTATION"):
    rect(slide, x, y, w, h, RGBColor(32, 39, 50), RGBColor(32, 39, 50))
    text(slide, title.title(), x+.22, y+.15, w-.44, .2, 9, RGBColor(168, 190, 220), False, "Arial")
    text(slide, value, x+.22, y+.45, w-.44, h-.58, 11.3, RGBColor(232, 236, 242), False, "Consolas")

def pill(slide, label, x, y, w, color=MINT):
    rect(slide, x, y, w, .38, PANEL2, color)
    text(slide, label, x, y+.08, w, .18, 10, color, True, "Arial", PP_ALIGN.CENTER)

def interview_slide(topic, questions, number):
    s = base(f"예상 질문 · {topic}", "Interview follow-up", number)
    text(s, "해당 페이지와 코드를 설명하다 보면 이러한 질문을 하실 수 있습니다.",
         .72, 1.38, 11.7, .34, 14, MUTED, False)
    colors = [MINT, BLUE, YELLOW]
    for i, (question, answer) in enumerate(questions):
        y = 1.92 + i * 1.55
        c = colors[i % len(colors)]
        text(s, f"Q{i+1}", .75, y+.08, .45, .25, 12, c, True, "Arial")
        text(s, question, 1.35, y, 3.45, .58, 15, WHITE, True)
        rect(s, 5.0, y-.08, 7.35, 1.22, PANEL, LINE)
        text(s, "답변", 5.28, y+.12, .65, .22, 10, c, True)
        text(s, answer, 6.0, y+.06, 6.02, .84, 12.5, WHITE)
    text(s, "답변 원칙  ·  현재 구현과 확장 계획을 구분하고, 측정하지 않은 성능은 수치로 주장하지 않습니다.",
         1.0, 6.7, 11.3, .28, 12, MUTED, True, align=PP_ALIGN.CENTER)
    return s

def engineering_detail_slide(title, premise, code_text, points, evidence, number):
    s = base(title, "Engineering detail", number)
    text(s, premise, .72, 1.38, 11.65, .48, 14, MUTED, False)
    code(s, code_text, .7, 1.98, 5.55, 4.45, "IMPLEMENTATION EVIDENCE")
    for i, (head, body, color) in enumerate(points):
        y = 2.0 + i * 1.02
        text(s, f"0{i+1}", 6.62, y+.12, .42, .22, 10, color, True, "Arial")
        text(s, head, 7.16, y+.05, 1.55, .28, 14, WHITE, True)
        text(s, body, 8.78, y+.02, 3.62, .55, 11.5, MUTED)
        rect(s, 6.62, y+.74, 5.72, .012, LINE, LINE, False)
    text(s, "근거", .72, 6.68, .55, .22, 11, BLUE, True)
    text(s, evidence, 1.36, 6.62, 10.95, .32, 11.5, WHITE, True)
    return s

# 1. Cover
s = prs.slides.add_slide(blank); s.background.fill.solid(); s.background.fill.fore_color.rgb = BG
add_image(s, GAME / "메인화면.png", 0, 0, 13.333, 7.5, .04, "fill")
overlay = rect(s, 0, 0, 13.333, 7.5, RGBColor(244,247,250), RGBColor(244,247,250), False); overlay.fill.transparency = 18
rect(s, .7, 1.0, .05, 4.6, BLUE, BLUE, False)
text(s, "LOCAL", 1.05, 1.05, 6.8, .72, 43, WHITE, True, "Arial")
text(s, "ARCADE", 1.05, 1.72, 6.8, .78, 48, MINT, True, "Arial")
text(s, "LAN 기반 실시간 멀티플레이 게임 플랫폼", 1.08, 2.72, 7.2, .4, 20, WHITE, True)
text(s, "게임 · 채팅 · AI 학습 퀴즈를 하나의 로컬 서비스로", 1.08, 3.25, 7.4, .34, 15, RGBColor(55, 65, 81), True)
text(s, "개인 프로젝트 · Full Stack · Backend 중심 설계", 1.08, 4.28, 6.8, .3, 13, WHITE, True)
text(s, "React + TypeScript  /  Spring Boot + MySQL  /  OpenAI Responses API", 1.08, 4.82, 7.7, .3, 11, MUTED, False, "Arial")
text(s, "2026  ·  기획 / 설계 / 구현 / 테스트", 1.08, 5.32, 5.5, .25, 10, MUTED, False, "Arial")

# 2. Agenda
s = base("Contents", "PORTFOLIO STRUCTURE", 2)
sections = [
    ("01", "Problem", "문제 정의 · 해결 방향 · 핵심 가설", MINT),
    ("02", "Opportunity", "타깃 · 경쟁 서비스 · 진입 전략 · 차별점", BLUE),
    ("03", "Product", "사용 흐름 · 주요 기능 · 실제 화면", PINK),
    ("04", "Engineering", "기술 선택 · 구조 · 동시성 · 트랜잭션 · AI", YELLOW),
    ("05", "Validation", "자동 테스트 · 측정 결과 · 한계와 확장", MINT),
    ("06", "Business & Me", "수익 가설 · 개인 기여 · 회고 · 다음 단계", BLUE),
]
for i, (n, head, body, c) in enumerate(sections):
    y = 1.38 + i * .82
    text(s, n, .75, y + .18, .55, .28, 13, c, True, "Arial")
    text(s, head, 1.55, y + .12, 2.05, .32, 17, WHITE, True, "Arial")
    text(s, body, 3.85, y + .15, 7.85, .32, 14, MUTED)
    rect(s, .75, y + .66, 11.55, .012, LINE, LINE, False)
text(s, "상세 페이지 구성", .75, 6.48, 1.55, .25, 11, BLUE, True)
text(s, "화면·코드로 먼저 이해 → 다음 장에서 설계 근거 → 예상 질문에 구현 범위로 답변", 2.35, 6.43, 9.85, .35, 13.5, WHITE, True)

# Pitch opening — establish the problem and opportunity before the technical proof.
problem_slide = base("왜 만들었는가: 함께 시작하기까지의 마찰", "01 · PROBLEM", 3)
problems = [
    ("START", "설치·계정·기기 준비", "짧게 함께 즐기려는 상황에서도 시작 전 준비가 경험보다 길어질 수 있습니다.", MINT),
    ("FRAGMENT", "게임과 소통의 분리", "게임, 공용 채팅, 점수와 운영 기능이 흩어지면 진행자가 여러 도구를 관리해야 합니다.", BLUE),
    ("TRUST", "공유 상태의 불일치", "여러 화면이 각자 결과를 만들면 턴·점수·보상에 대한 하나의 기준을 유지하기 어렵습니다.", YELLOW),
]
for i, (tag, head, body, c) in enumerate(problems):
    x = .72 + i * 4.08
    rect(problem_slide, x, 1.62, 3.7, 3.58, PANEL, LINE)
    pill(problem_slide, tag, x+.28, 1.92, 1.0, c)
    text(problem_slide, head, x+.28, 2.55, 3.1, .58, 19, WHITE, True)
    text(problem_slide, body, x+.28, 3.36, 3.08, 1.15, 13, MUTED)
text(problem_slide, "프로젝트 가설", .75, 5.72, 1.25, .25, 11, MINT, True)
text(problem_slide, "같은 공간의 소규모 사용자는 ‘빠른 접속’과 ‘서버가 보장하는 동일한 결과’를 함께 원한다.", 2.05, 5.63, 10.1, .48, 17, WHITE, True)
text(problem_slide, "※ 시장 전체의 불편을 단정한 것이 아니라, 교실·친구 모임에서 검증할 제품 가설로 정의했습니다.", .75, 6.5, 11.55, .3, 11.5, MUTED, align=PP_ALIGN.CENTER)
problem_id = prs.slides._sldIdLst[-1]

solution_slide = base("해결 방향: 코드 하나로 연결되는 로컬 서비스", "01 · SOLUTION", 4)
solution_steps = [
    ("01", "즉시 접속", "서버가 발급한 6자리 코드와 닉네임만으로 동일 LAN에서 참여", MINT),
    ("02", "하나의 경험", "게임·공용 채팅·크레딧·랭킹·AI 퀴즈·관리 기능 통합", BLUE),
    ("03", "서버 판정", "턴·결과·참가비·보상을 서버가 검증해 모든 화면에 같은 상태 제공", YELLOW),
    ("04", "반복 검증", "최대 6명 동시 요청을 자동화하고 성공률·중복 차감·응답 시간을 기록", PINK),
]
for i, (n, head, body, c) in enumerate(solution_steps):
    y = 1.5 + i * 1.18
    text(solution_slide, n, .78, y+.18, .55, .28, 13, c, True, "Arial")
    text(solution_slide, head, 1.55, y+.12, 2.05, .34, 17, WHITE, True)
    rect(solution_slide, 3.82, y-.02, 8.5, .82, PANEL, LINE)
    text(solution_slide, body, 4.12, y+.18, 7.9, .38, 13.5, MUTED)
text(solution_slide, "한 줄 정의", .78, 6.5, 1.05, .24, 11, BLUE, True)
text(solution_slide, "설치 부담 없이 접속하고, 서버가 동일한 게임 상태를 책임지는 LAN 기반 멀티플레이 플랫폼", 1.9, 6.43, 10.3, .38, 15, WHITE, True)
solution_id = prs.slides._sldIdLst[-1]

market_slide = base("누구에게 필요한가: 작고 가까운 그룹부터", "02 · TARGET & OPPORTUNITY", 5)
targets = [
    ("교실 · 실습실", "수업 전후 짧은 활동과 AI 학습 퀴즈", "관리자가 문제 생성과 점수를 운영", MINT),
    ("친구 · 가족 모임", "한 화면을 함께 보며 바로 시작하는 미니게임", "개별 브라우저는 입력과 상태 확인에 사용", BLUE),
    ("동아리 · 워크숍", "6명 이하 소그룹의 아이스브레이킹", "테마형 콘텐츠와 결과 화면으로 반복 운영", YELLOW),
]
for i, (head, use, value, c) in enumerate(targets):
    x = .72 + i * 4.08
    rect(market_slide, x, 1.55, 3.72, 3.72, PANEL, LINE)
    text(market_slide, f"0{i+1}", x+.3, 1.88, .5, .25, 12, c, True, "Arial")
    text(market_slide, head, x+.3, 2.3, 3.05, .45, 18, WHITE, True)
    text(market_slide, "사용 장면", x+.3, 3.02, .9, .22, 10, c, True)
    text(market_slide, use, x+.3, 3.35, 3.02, .66, 12.5, MUTED)
    text(market_slide, "제공 가치", x+.3, 4.22, .9, .22, 10, c, True)
    text(market_slide, value, x+.3, 4.53, 3.02, .58, 12.5, WHITE, True)
text(market_slide, "초기 검증 지표", .75, 5.75, 1.35, .24, 11, BLUE, True)
text(market_slide, "입장 완료율 · 첫 게임 시작 시간 · 세션당 재경기 수 · AI 퀴즈 완료율 · 운영자 개입 횟수", 2.2, 5.68, 9.95, .38, 14, WHITE, True)
text(market_slide, "시장 규모 수치보다 실제 소그룹 파일럿에서 반복 사용 여부를 먼저 확인합니다.", .75, 6.5, 11.55, .3, 12, MUTED, align=PP_ALIGN.CENTER)
market_id = prs.slides._sldIdLst[-1]

competition_slide = base("경쟁 서비스와의 포지션", "02 · COMPETITION", 6)
cols = [(.72, 2.15, "서비스"), (2.92, 2.25, "시작 방식"), (5.24, 2.15, "핵심 경험"), (7.49, 2.1, "공유 상태"), (9.68, 2.85, "Local Arcade와의 차이")]
for x, w, label in cols:
    rect(competition_slide, x, 1.48, w, .62, PANEL2, LINE)
    text(competition_slide, label, x+.08, 1.68, w-.16, .22, 11, WHITE, True, align=PP_ALIGN.CENTER)
rows = [
    ("Jackbox", "호스트 구매·공유 화면\n참가자는 웹 기기", "파티형 게임", "호스트 화면 중심", "구매형 콘텐츠와 원격 화면 공유에 강점"),
    ("Gartic Phone", "닉네임·초대 링크\n음성 통화 권장", "그리기·문장 전달", "라운드 결과 공유", "단일 놀이 흐름에 집중"),
    ("Board Game Arena", "브라우저 계정·온라인", "보드게임 플레이", "서버 기반 게임 상태", "방대한 정식 보드게임 카탈로그"),
    ("Local Arcade", "LAN·6자리 코드\n닉네임", "미니게임+채팅+AI 퀴즈", "서버 판정+공유 크레딧", "소규모 현장 운영과 기술 검증을 한 서비스에 통합"),
]
for i, row in enumerate(rows):
    y = 2.13 + i * 1.0
    fill = PANEL2 if i == 3 else PANEL
    color = MINT if i == 3 else WHITE
    widths = [2.15, 2.25, 2.15, 2.1, 2.85]
    xs = [.72, 2.92, 5.24, 7.49, 9.68]
    for j, value in enumerate(row):
        rect(competition_slide, xs[j], y, widths[j], .94, fill, LINE)
        text(competition_slide, value, xs[j]+.1, y+.14, widths[j]-.2, .62, 10.5 if j else 12, color if j == 0 else MUTED, j == 0, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
text(competition_slide, "비교 기준은 기능의 우열이 아니라 시작 환경과 서비스 범위입니다.", .75, 6.32, 11.5, .28, 12, WHITE, True, align=PP_ALIGN.CENTER)
text(competition_slide, "Sources · jackboxgames.com/how-to-play · garticphone.com · en.doc.boardgamearena.com  (2026-08-11 확인)", .75, 6.72, 11.5, .22, 9.5, MUTED, align=PP_ALIGN.CENTER)
competition_id = prs.slides._sldIdLst[-1]

entry_slide = base("시장 진입 가능성: 작은 파일럿에서 반복 사용으로", "02 · GO-TO-MARKET", 7)
stages = [
    ("01", "Portfolio Demo", "GitHub·발표·면접에서 실행 가능한 데모와 테스트 보고서 공개", "관찰: 입장 성공·게임 완료", MINT),
    ("02", "Closed Pilot", "교실·친구 모임의 4~6명 세션에서 진행자 없이 시작 가능한지 확인", "검증: 재경기·이탈·오류", BLUE),
    ("03", "Theme Package", "행사·수업 목적에 맞춘 퀴즈·배너·게임 구성을 패키지화", "검증: 재사용·운영 시간", YELLOW),
    ("04", "Scale Decision", "원격 접속 수요가 확인되면 WebSocket·Room 분리·Redis를 단계 도입", "판단: 비용 대비 수요", PINK),
]
for i, (n, head, body, metric, c) in enumerate(stages):
    x = .65 + i * 3.12
    rect(entry_slide, x, 1.58, 2.78, 4.72, PANEL, LINE)
    text(entry_slide, n, x+.25, 1.9, .45, .25, 12, c, True, "Arial")
    text(entry_slide, head, x+.25, 2.38, 2.28, .55, 17, WHITE, True)
    text(entry_slide, body, x+.25, 3.16, 2.28, 1.42, 12.2, MUTED)
    rect(entry_slide, x+.2, 4.88, 2.38, .9, PANEL2, c)
    text(entry_slide, metric, x+.34, 5.12, 2.1, .4, 11, c, True, align=PP_ALIGN.CENTER)
    if i < 3:
        text(entry_slide, "→", x+2.76, 3.65, .36, .4, 20, MUTED, True, "Arial", PP_ALIGN.CENTER)
text(entry_slide, "대규모 배포를 먼저 가정하지 않고, 현재 6명 요구에서 얻은 근거로 다음 투자를 결정합니다.", .75, 6.64, 11.55, .3, 12.5, WHITE, True, align=PP_ALIGN.CENTER)
entry_id = prs.slides._sldIdLst[-1]

differentiation_slide = base("차별점: 기능 수보다 연결 방식과 검증 근거", "02 · DIFFERENTIATION", 8)
diffs = [
    ("LAN-FIRST", "설치·계정 없이 6자리 코드로 참여", "동일 공간의 최대 6명이라는 명확한 운영 범위", MINT),
    ("ONE STATE", "게임·채팅·크레딧·랭킹을 한 서버 상태로 연결", "클라이언트가 아닌 서버가 최종 판정", BLUE),
    ("MEASURED", "6명 동시 입장·참가·중복 거절·60회 조회 자동화", "Python 코드와 Markdown 리포트로 재현", YELLOW),
    ("GROUNDED AI", "로컬 위키 발췌→Schema→서버 검증→DB 저장", "외부 API 장애 시 저장 문제로 대체", PINK),
]
for i, (tag, head, proof, c) in enumerate(diffs):
    x = .72 + (i % 2) * 6.12
    y = 1.55 + (i // 2) * 2.42
    rect(differentiation_slide, x, y, 5.72, 2.08, PANEL, LINE)
    pill(differentiation_slide, tag, x+.28, y+.26, 1.35, c)
    text(differentiation_slide, head, x+.28, y+.88, 5.05, .48, 15, WHITE, True)
    text(differentiation_slide, proof, x+.28, y+1.47, 5.05, .32, 11.5, MUTED)
text(differentiation_slide, "차별점의 각 문장은 뒤의 화면·코드·테스트 결과 페이지에서 근거를 확인할 수 있습니다.", .75, 6.56, 11.55, .3, 13, WHITE, True, align=PP_ALIGN.CENTER)
differentiation_id = prs.slides._sldIdLst[-1]

profile_slide = base("개인 프로젝트 · 역할 개요", "02 · ABOUT THE BUILDER", 9)
add_image(profile_slide, ROOT / "Yuni.png", .72, 1.52, 3.72, 4.62, mode="contain")
text(profile_slide, "Full Stack Developer", 4.88, 1.6, 5.7, .38, 20, WHITE, True, "Arial")
text(profile_slide, "기획부터 검증 자료 제작까지 1인 수행", 4.88, 2.08, 6.55, .38, 15, MUTED)
roles = [
    ("Product", "사용 흐름·게임 규칙·관리자 운영·프로모션 정책", MINT),
    ("Frontend", "React 화면·상태 표현·애니메이션·반응형 UI", BLUE),
    ("Backend", "Spring API·게임 판정·명령 큐·세션 정리", YELLOW),
    ("Data & AI", "MySQL 트랜잭션·위키 기반 문제 생성·검증·저장", PINK),
    ("Validation", "JUnit 규칙 테스트·Python 동시 요청·결과 리포트", MINT),
]
for i, (head, body, c) in enumerate(roles):
    y = 2.78 + i * .68
    text(profile_slide, head, 4.9, y+.06, 1.32, .25, 12, c, True, "Arial")
    text(profile_slide, body, 6.4, y, 5.75, .38, 12.5, WHITE)
text(profile_slide, "※ 실제 인물 사진과 이름은 공개 범위에 맞춰 최종 제출본에서 교체할 수 있도록 독립된 프로필 영역으로 구성", .75, 6.64, 11.55, .26, 10.5, MUTED, align=PP_ALIGN.CENTER)
profile_id = prs.slides._sldIdLst[-1]

# 3. Overview
s = base("프로젝트 소개", "01 · OVERVIEW", 3)
add_image(s, GAME / "참여코드 입력 화면.png", .65, 1.48, 5.7, 2.65)
add_image(s, GAME / "게임 종류 및 메뉴 창.png", .65, 4.38, 5.7, 1.42)
rect(s, 6.7, 1.48, 5.95, 4.95, PANEL, LINE)
text(s, "참여 코드 하나로 시작하는 로컬 아케이드", 7.05, 1.8, 5.2, .55, 22, WHITE, True)
rich_lines(s, [("형태  ", "개인 프로젝트 · 기여도 100%", MINT),("대상  ", "동일 LAN 내 최대 6명", PINK),("접속  ", "6자리 코드 + 닉네임", YELLOW),("콘텐츠  ", "야추 · 레이싱 · 행맨 · 마피아 · AI 퀴즈", BLUE),("운영  ", "공용 채팅 · 크레딧 · 랭킹 · 관리자 페이지", MINT)], 7.05, 2.65, 5.05, 2.75, 16)
text(s, "목표: 화면 구현을 넘어 여러 사용자가 공유하는 상태를 서버 중심으로 일관되게 관리", 7.05, 5.55, 5.0, .58, 14, WHITE, True)

# 4. Stack — explain responsibility and selection instead of listing logos.
s = base("기술 구성과 역할", "02 · Technology", 4)
stack_rows = [
    ("Client", "React 18 · TypeScript · Vite", "사용자 입력, 게임 화면과 애니메이션. 서버 상태를 표현하는 역할로 제한", MINT),
    ("Server", "Java 21 · Spring Boot 3.4", "게임 규칙 검증, 공유 상태 변경, 명령 순서와 세션 관리", BLUE),
    ("Data", "MySQL 8 · Spring Data JPA", "플레이어 크레딧과 검증된 퀴즈를 영속화하고 트랜잭션 경계 제공", YELLOW),
    ("AI", "OpenAI Responses API · JSON Schema", "로컬 위키 기반 문제 생성 후 형식·정답·출처를 서버에서 재검증", PINK),
    ("Network", "HTTP REST · Bearer Token · LAN", "1초 폴링으로 상태 동기화. 현재 규모에서 단순성과 복구 용이성 선택", MINT),
]
for i, (area, tech, reason, c) in enumerate(stack_rows):
    y = 1.48 + i * 1.03
    text(s, area, .75, y + .16, 1.2, .28, 13, c, True, "Arial")
    text(s, tech, 2.05, y + .12, 3.35, .32, 15, WHITE, True)
    text(s, reason, 5.65, y + .10, 6.45, .48, 12.5, MUTED)
    rect(s, .75, y + .76, 11.5, .012, LINE, LINE, False)

# 4. Architecture
s = base("System Architecture", "03 · DESIGN", 4)
text(s, "브라우저 1~6명", .65, 1.55, 2.2, .35, 16, WHITE, True)
for i in range(3):
    rect(s,.7+i*.72,2.0,.55,.72,PANEL2,MINT); text(s,"USER",.7+i*.72,2.25,.55,.15,7,MINT,True,"Arial",PP_ALIGN.CENTER)
text(s,"HTTP · 1초 폴링",3.0,2.17,1.55,.25,11,MUTED,True,"Arial",PP_ALIGN.CENTER)
text(s,"→",4.35,2.02,.5,.5,27,MINT,True,"Arial",PP_ALIGN.CENTER)
rect(s,4.9,1.63,2.2,1.35,PANEL,MINT); text(s,"React + TypeScript",5.05,1.98,1.9,.3,15,WHITE,True,"Arial",PP_ALIGN.CENTER)
text(s,"REST API · Token",7.2,2.17,1.55,.25,11,MUTED,True,"Arial",PP_ALIGN.CENTER); text(s,"→",8.55,2.02,.5,.5,27,PINK,True,"Arial",PP_ALIGN.CENTER)
rect(s,9.05,1.63,3.05,1.35,PANEL,PINK); text(s,"Spring Boot Server",9.25,1.98,2.65,.3,16,WHITE,True,"Arial",PP_ALIGN.CENTER)
services=[("GAME STATE","턴 · 결과 · 크레딧",MINT),("COMMAND QUEUE","요청 순차 처리",PINK),("QUIZ PIPELINE","생성 · 검증 · 저장",YELLOW)]
for i,(a,b,c) in enumerate(services):
    x=.8+i*4.15; rect(s,x,3.65,3.7,1.1,PANEL,LINE); text(s,a,x+.2,3.91,3.3,.22,12,c,True,"Arial",PP_ALIGN.CENTER); text(s,b,x+.2,4.3,3.3,.2,12,WHITE,False,"Arial",PP_ALIGN.CENTER)
text(s,"↓",3.0,5.0,.4,.4,25,MINT,True,"Arial",PP_ALIGN.CENTER); text(s,"↓",9.9,5.0,.4,.4,25,YELLOW,True,"Arial",PP_ALIGN.CENTER)
rect(s,1.15,5.55,4.2,.75,PANEL2,BLUE); text(s,"MySQL · 플레이어 / 문제 / 크레딧",1.35,5.79,3.8,.25,13,WHITE,True,"Arial",PP_ALIGN.CENTER)
rect(s,7.6,5.55,4.2,.75,PANEL2,YELLOW); text(s,"LLM Wiki → OpenAI Responses API",7.8,5.79,3.8,.25,13,WHITE,True,"Arial",PP_ALIGN.CENTER)
text(s,"최종 판정은 서버가 담당하고 클라이언트는 입력과 표현을 담당",3.2,6.65,6.9,.28,13,MUTED,True,"Arial",PP_ALIGN.CENTER)

# 5. User flow
s = base("User Flow", "04 · EXPERIENCE", 5)
steps=[("01","서버 실행","참여 코드 발급"),("02","빠른 입장","코드 + 닉네임"),("03","로비","게임 선택"),("04","멀티플레이","참가 · 턴 · 채팅"),("05","결과 정산","크레딧 · 랭킹")]
for i,(n,a,b) in enumerate(steps):
    x=.6+i*2.52; rect(s,x,1.65,2.2,1.25,PANEL, [MINT,PINK,YELLOW,BLUE,MINT][i]); text(s,n,x+.15,1.82,.45,.25,12,[MINT,PINK,YELLOW,BLUE,MINT][i],True,"Arial"); text(s,a,x+.15,2.15,1.9,.25,15,WHITE,True); text(s,b,x+.15,2.51,1.9,.2,11,MUTED)
    if i<4: text(s,"→",x+2.2,2.0,.32,.4,20,MUTED,True,"Arial",PP_ALIGN.CENTER)
add_image(s, GAME / "메인화면.png", .65, 3.35, 7.3, 3.25)
add_image(s, GAME / "관리자 페이지 플레이어 점수 관리 및 AI문제 출제 활성화 버튼.png", 8.25, 3.35, 4.4, 3.25)

# 6. Yacht
s=base("턴과 상태를 서버가 관리하는 멀티플레이 야추", "05 · CORE GAME", 6)
add_image(s,GAME/"야추게임화면.png",.65,1.45,5.25,5.45)
rich_lines(s,[("READY  ","2명 이상, 동일 참가 점수 검증",MINT),("ROLL  ","턴당 최대 3회 · HOLD 유지",PINK),("SCORE  ","12개 항목과 상단 보너스 계산",YELLOW),("RESULT  ","공동 우승 · 기권 · 환불 · 재경기",BLUE)],6.25,1.62,5.8,2.1,15)
code(s,"YachtGame g = activeTurn(player);\n\nif (!visibleDice.equals(serverDice))\n    throw new IllegalArgumentException(...);\n\nif (sheet.containsKey(category))\n    throw new IllegalArgumentException(...);\n\nint points = calculate(category, g.dice);\nsheet.put(category, points);",6.25,3.85,6.0,2.78,"SERVER-SIDE VALIDATION")

# 7. Racing
s=base("모든 플레이어가 같은 결과를 공유하는 레이싱", "06 · SYNCHRONIZATION", 7)
add_image(s,GAME/"달팽이 레이싱 .png",.65,1.5,5.8,5.35)
add_image(s,GAME/"달팽이 레이싱 종료 화면.png",6.7,1.5,2.8,2.45)
code(s,"List<Integer> order = shuffle(1, 2, 3, 4, 5);\nlong raceId = System.currentTimeMillis();\n\n// 서버가 단일 결과를 생성\nlastResult = new RaceResult(raceId, order);\n\n// 브라우저는 전달받은 순위만 표현\nsetRaceOrder(out.state.race.lastResult.order);",9.75,1.5,2.9,2.45,"ONE SOURCE OF TRUTH")
rich_lines(s,[("① 판정  ","랜덤 순위와 보상은 서버에서 한 번만",MINT),("② 식별  ","raceId로 새 경기와 중복 실행 구분",PINK),("③ 표현  ","카운트다운·주행·시상대는 클라이언트",YELLOW)],6.7,4.35,5.75,1.55,14)
text(s,"사용자마다 자체 랜덤 결과를 만들지 않아 모든 화면의 우승자가 일치",6.7,6.25,5.65,.4,14,WHITE,True)

# 8. AI quiz
s=base("근거 기반 AI 퀴즈 파이프라인", "07 · GENERATIVE AI", 8)
add_image(s,GAME/"LLM WIKI  를기반으로 OPenAI 가문제를 출제 하는 화면.png",.6,1.5,4.2,4.9)
flow=[("WIKI","Markdown 발췌",MINT),("GENERATE","Responses API",PINK),("STRUCTURE","JSON Schema",YELLOW),("VALIDATE","출처·정답 검증",BLUE),("STORE","MySQL 재사용",MINT)]
for i,(a,b,c) in enumerate(flow):
    x=5.05+i*1.48; rect(s,x,1.55,1.28,1.12,PANEL,c); text(s,a,x+.08,1.82,1.12,.18,10,c,True,"Arial",PP_ALIGN.CENTER); text(s,b,x+.08,2.16,1.12,.24,9,WHITE,False,"Arial",PP_ALIGN.CENTER)
    if i<4:text(s,"→",x+1.25,1.9,.25,.3,14,MUTED,True,"Arial",PP_ALIGN.CENTER)
code(s,'아래 LLM Wiki 발췌문만 근거로\n한국어 프로그래밍 퀴즈 한 문제를 만들어라.\n\n"choices": { "minItems": 2, "maxItems": 4 },\n"answer":  { "minimum": 0, "maximum": 3 }\n\nif (excerpts.noneMatch(q.source()))\n    throw new IllegalArgumentException("출처 불일치");',5.05,3.02,7.35,3.38,"PROMPT → SCHEMA → SERVER VALIDATION")
text(s,"AI 생성 OFF 시 검증된 저장 문제로 자동 대체",8.0,6.57,4.4,.28,12,YELLOW,True,"Arial",PP_ALIGN.RIGHT)

# 9. Other games
s=base("게임 규칙을 상태 머신으로 확장", "08 · GAME CONTENT", 9)
add_image(s,GAME/"낮시 간채팅으로 소통.png",.65,1.5,3.9,4.35)
add_image(s,GAME/"행맨 진행.png",4.75,1.5,3.9,4.35)
add_image(s,GAME/"문제 정답 화면 .png",8.85,1.5,3.8,4.35)
text(s,"마피아",.65,6.02,3.9,.25,15,PINK,True,align=PP_ALIGN.CENTER); text(s,"게임 상태에 종속된 낮/밤 전용 채팅",.65,6.38,3.9,.25,11,MUTED,align=PP_ALIGN.CENTER)
text(s,"1:1 행맨",4.75,6.02,3.9,.25,15,YELLOW,True,align=PP_ALIGN.CENTER); text(s,"라운드 · 피해량 · 승패 판정",4.75,6.38,3.9,.25,11,MUTED,align=PP_ALIGN.CENTER)
text(s,"위키 퀴즈",8.85,6.02,3.8,.25,15,MINT,True,align=PP_ALIGN.CENTER); text(s,"문제 · 해설 · 출처 · 보상",8.85,6.38,3.8,.25,11,MUTED,align=PP_ALIGN.CENTER)

# 10. Character
s=base("시간 기반 2D 캐릭터 컨트롤", "09 · ANIMATION", 10)
add_image(s,GAME/"캐릭터 조작 테스트.png",.65,1.5,6.15,4.95)
code(s,"const dt = Math.min((now - previous) / 1000, 0.032);\nconst velocity = running ? RUN_SPEED : WALK_SPEED;\n\nx += direction * velocity * dt;\nvy -= GRAVITY * dt;\ny  += vy * dt;\n\nif (landingOnPlatform) {\n  y = platformTop; vy = 0; grounded = true;\n}\nraf = requestAnimationFrame(tick);",7.15,1.5,5.5,3.85,"REQUESTANIMATIONFRAME + ΔT")
rich_lines(s,[("INPUT  ","방향키 · Shift 달리기 · 점프",MINT),("PHYSICS  ","중력 · 수직 속도 · 발판 착지",PINK),("STATE  ","idle / walk / run / jump 전환",YELLOW)],7.15,5.65,5.2,1.1,13)

# 11. Social/admin
s=base("공용 채팅과 크레딧 운영 시스템", "10 · OPERATIONS", 11)
add_image(s,GAME/"채팅화면.png",.65,1.48,3.65,4.7)
add_image(s,GAME/"관리자 페이지 플레이어 점수 관리 및 AI문제 출제 활성화 버튼.png",4.55,1.48,4.55,4.7)
code(s,"String clean = text.strip();\nif (clean.isBlank() || clean.length() > 200)\n    throw new IllegalArgumentException(...);\n\nchat.addLast(message);\nwhile (chat.size() > 80) chat.removeFirst();\n\n@Transactional\nvoid grantCredits(Player p, long amount) { ... }",9.35,1.48,3.3,3.42,"VALIDATE · LIMIT · TRANSACTION")
rich_lines(s,[("공용 채팅  ","게임 종류와 무관한 최근 80개 메시지",MINT),("마피아 채팅  ","낮/밤·역할 상태에 따라 별도 접근 제어",BLUE),("관리자  ","코드 인증 · 지급/차감 · 퇴장",PINK)],9.35,5.02,3.2,1.72,11.5)

# 12. Reliability
s=base("서버 중심 상태 관리", "11 · RELIABILITY", 12)
problems=[("동시 요청","GameCommandQueue + synchronized","턴과 공유 상태 변경을 순서대로 처리",PINK),("데이터 일관성","@Transactional","참가비 차감과 보상 지급을 원자적으로 처리",MINT),("클라이언트 조작","서버 재검증","턴·주사위·참가 금액·관리자 요청 검증",YELLOW),("유령 참가자","활동 시간 정리","비활성 사용자를 주기적으로 제거",BLUE)]
for i,(p,sol,desc,c) in enumerate(problems):
    y=1.48+i*1.28; rect(s,.75,y,2.35,.92,PANEL,c); text(s,p,.95,y+.28,1.95,.25,15,c,True,align=PP_ALIGN.CENTER)
    text(s,"→",3.25,y+.2,.5,.4,24,MUTED,True,"Arial",PP_ALIGN.CENTER)
    rect(s,3.85,y,3.1,.92,PANEL2,LINE); text(s,sol,4.05,y+.25,2.7,.28,14,WHITE,True,"Consolas",PP_ALIGN.CENTER)
    rect(s,7.2,y,5.25,.92,PANEL,LINE); text(s,desc,7.48,y+.25,4.7,.32,13,WHITE,False)
text(s,"핵심 원칙  ·  화면의 상태보다 서버의 상태를 신뢰한다",3.15,6.72,7.0,.3,15,MINT,True,"Arial",PP_ALIGN.CENTER)

# 13. Testing
s=base("검증 가능한 게임 규칙", "12 · TESTING", 13)
code(s,'@Test\nvoid yachtScoresAreCalculatedByRule() {\n  assertEquals(50, calculate("YACHT", 6,6,6,6,6));\n  assertEquals(30, calculate("LARGE_STRAIGHT", 1,2,3,4,5));\n}\n\n@Test\nvoid mafiaNightAndVoteFlowIsDeterministic() {\n  // 역할, 제거, 승리 조건을 단계별 검증\n}',.7,1.55,5.65,4.9,"JUNIT · DOMAIN RULE TEST")
rect(s,6.7,1.55,5.9,4.9,PANEL,LINE)
text(s,"테스트 대상",7.05,1.9,2.2,.3,17,MINT,True)
rich_lines(s,[("야추 점수  ","12개 카테고리와 스트레이트·보너스",MINT),("마피아 흐름  ","역할 배정·낮/밤·투표·승리 조건",PINK),("서버 검증  ","잘못된 턴과 중복 요청 거부",YELLOW)],7.05,2.45,5.0,2.0,15)
text(s,"확장 계획",7.05,4.9,2.2,.3,17,BLUE,True)
text(s,"Controller 통합 테스트 · 동시성 테스트 · 프론트 UI 테스트",7.05,5.4,4.9,.7,14,WHITE,True)

# 14. Learning
s=base("Troubleshooting & Learning", "13 · RETROSPECTIVE", 14)
items=[("사용자별 결과 불일치","랜덤 판정을 서버에서 한 번만 수행",MINT),("턴·크레딧 충돌","명령 큐 + 동기화 + 트랜잭션",PINK),("LLM 형식·근거 오류","근거 제한 + JSON Schema + 서버 검증",YELLOW),("환경별 이동 속도 차이","프레임 시간 차이 dt 적용",BLUE)]
for i,(a,b,c) in enumerate(items):
    x=.7+(i%2)*6.15; y=1.55+(i//2)*2.25
    rect(s,x,y,5.8,1.75,PANEL,c); text(s,f"0{i+1}",x+.25,y+.25,.55,.3,14,c,True,"Arial"); text(s,a,x+.95,y+.22,4.45,.35,17,WHITE,True); text(s,b,x+.95,y+.83,4.45,.45,14,MUTED)
rect(s,1.55,6.15,10.2,.62,PANEL2,MINT); text(s,"기능을 만드는 것에서 끝내지 않고, 신뢰 가능한 서비스 상태를 설계하는 경험",1.75,6.35,9.8,.25,15,WHITE,True,align=PP_ALIGN.CENTER)
retrospective_id = prs.slides._sldIdLst[-1]

# 15. Result
s=base("Result", "22 · OUTCOME", 23)
imgs=["야추게임화면.png","달팽이 레이싱 종료 화면.png","문제 정답 화면 .png","캐릭터 조작 테스트.png","관리자 페이지 플레이어 점수 관리 및 AI문제 출제 활성화 버튼.png"]
for i,name in enumerate(imgs):
    add_image(s,GAME/name,.6+i*2.48,1.42,2.25,2.48)
metrics=[("6","MAX PLAYERS"),("5+","GAME CONTENTS"),("1","SHARED CREDIT"),("100%","PERSONAL BUILD")]
for i,(n,lbl) in enumerate(metrics):
    x=.75+i*3.05; rect(s,x,4.35,2.7,1.25,PANEL,LINE); text(s,n,x+.15,4.57,2.4,.42,25,[MINT,PINK,YELLOW,BLUE][i],True,"Arial",PP_ALIGN.CENTER); text(s,lbl,x+.15,5.12,2.4,.2,9,MUTED,True,"Arial",PP_ALIGN.CENTER)
text(s,"게임 · 실시간 상태 · AI · 운영 기능을 하나의 풀스택 서비스로 통합",1.4,6.15,10.5,.42,19,WHITE,True,align=PP_ALIGN.CENTER)
text(s,"THANK YOU",4.8,6.73,3.7,.3,15,MINT,True,"Arial",PP_ALIGN.CENTER)
result_slide = s
result_id = prs.slides._sldIdLst[-1]

# 16. Request traffic model
s=base("반복 요청이 만드는 서버 부하", "DEEP DIVE · REQUEST TRAFFIC", 15)
rect(s,.65,1.48,4.0,4.9,PANEL,LINE)
text(s,"현재 통신 모델",.95,1.82,3.3,.35,19,MINT,True)
rich_lines(s,[("게임 상태  ","/api/state · 1초마다",MINT),("소셜 상태  ","/api/social/state · 1초마다",PINK),("마피아 상태  ","참여 중 추가 조회",YELLOW),("관리자  ","2초마다 2개 운영 API 조회",BLUE)],.95,2.4,3.3,2.3,14)
text(s,"최대 6명 기준",.95,5.12,3.2,.3,14,MUTED,True)
text(s,"기본 약 12 read req/s",.95,5.55,3.2,.4,21,WHITE,True,"Arial")
text(s,"※ 계산값이며 부하 테스트 측정치는 아님",.95,6.02,3.2,.2,9,MUTED)
flow=[("READ","상태 조회","병렬 처리",MINT),("COMMAND","게임 행동","큐 진입",PINK),("DB WRITE","크레딧·문제","트랜잭션",YELLOW)]
for i,(a,b,detail,c) in enumerate(flow):
    x=5.15+i*2.5; rect(s,x,1.7,2.1,1.28,PANEL,c); text(s,a,x+.12,1.98,1.86,.22,12,c,True,"Arial",PP_ALIGN.CENTER); text(s,b,x+.12,2.4,1.86,.2,11,WHITE,False,"Arial",PP_ALIGN.CENTER)
    if i<2:text(s,"→",x+2.1,2.1,.4,.3,17,MUTED,True,"Arial",PP_ALIGN.CENTER)
code(s,"// 상태 조회는 큐를 거치지 않음\n@GetMapping(\"/state\")\nObject state(...) { return service.state(player); }\n\n// 상태 변경 명령만 순차 처리\n@PostMapping(\"/yacht/score\")\nObject score(...) {\n  return queue.run(() -> social.score(...));\n}",5.15,3.4,7.5,2.98,"READ / WRITE PATH SEPARATION")
text(s,"핵심: 모든 요청을 직렬화하지 않고, 충돌 위험이 큰 상태 변경 명령을 통제",5.15,6.55,7.45,.3,13,WHITE,True)

# 17. Queue
s=base("명령 큐로 순서와 과부하를 통제", "DEEP DIVE · COMMAND QUEUE", 16)
code(s,"new ThreadPoolExecutor(\n  1, 1, 0L, MILLISECONDS,\n  new ArrayBlockingQueue<>(128),\n  threadFactory,\n  new AbortPolicy()\n);\n\nreturn CompletableFuture\n  .supplyAsync(command, executor).join();",.7,1.5,5.3,4.85,"GAMECOMMANDQUEUE.JAVA")
steps=[("①","단일 Worker","게임 명령 실행 순서 보장",MINT),("②","Bounded 128","무제한 적재로 인한 메모리 증가 방지",PINK),("③","Abort Policy","포화 시 즉시 거절하고 재시도 안내",YELLOW),("④","Exception unwrap","도메인 오류를 원래 메시지로 반환",BLUE)]
for i,(n,a,b,c) in enumerate(steps):
    y=1.5+i*1.22; rect(s,6.35,y,6.25,.94,PANEL,LINE); text(s,n,6.55,y+.27,.45,.25,14,c,True,"Arial"); text(s,a,7.1,y+.21,1.65,.25,14,WHITE,True); text(s,b,8.75,y+.22,3.55,.32,12,MUTED)
rect(s,6.35,6.42,6.25,.48,PANEL2,PINK); text(s,"Trade-off · 처리량보다 작은 LAN 환경의 상태 일관성을 우선",6.55,6.57,5.85,.2,11,WHITE,True,align=PP_ALIGN.CENTER)

# 18. Transactions
s=base("트랜잭션 경계를 게임 결과와 맞추다", "DEEP DIVE · TRANSACTION", 17)
text(s,"예: 야추 점수 확정과 크레딧 정산",.75,1.55,4.3,.35,18,YELLOW,True)
timeline=[("검증","현재 턴·주사위·카테고리"),("계산","서버 주사위로 점수 산출"),("상태","점수표·다음 턴 갱신"),("정산","우승자 크레딧 저장")]
for i,(a,b) in enumerate(timeline):
    x=.7+i*3.05; c=[MINT,PINK,YELLOW,BLUE][i]; rect(s,x,2.15,2.55,1.25,PANEL,c); text(s,a,x+.15,2.43,2.25,.25,15,c,True,align=PP_ALIGN.CENTER); text(s,b,x+.15,2.82,2.25,.3,11,WHITE,align=PP_ALIGN.CENTER)
    if i<3:text(s,"→",x+2.55,2.52,.5,.35,20,MUTED,True,"Arial",PP_ALIGN.CENTER)
code(s,"@Transactional\npublic synchronized Map<String,Object> score(...) {\n  YachtGame g = activeTurn(player);\n  validateDiceAndCategory(g, request);\n  sheet.put(category, calculate(category, g.dice));\n  if (finished) finishYacht(g); // repo.save(winner)\n  else advanceTurn(g);\n  return state(player);\n}",.7,3.85,6.0,2.72,"APPLICATION TRANSACTION BOUNDARY")
rich_lines(s,[("원자성  ","예외 발생 시 DB 변경 롤백",MINT),("일관성  ","허용 금액과 잔액을 도메인에서 검증",PINK),("격리 보완  ","명령 큐 + synchronized로 인메모리 경쟁 제어",YELLOW),("주의  ","인메모리 상태는 DB 롤백 대상이 아니므로 확장 시 영속 상태화 필요",BLUE)],7.05,3.92,5.25,2.55,13)

# 19. Failure scenario
s=base("실패 시나리오로 보는 데이터 일관성", "DEEP DIVE · FAILURE HANDLING", 18)
scenarios=[("동시 참가","두 사용자가 마지막 크레딧으로 동시에 참가","큐가 명령을 순차 실행 → 두 번째 요청은 최신 잔액으로 재검증",MINT),("중복 점수","같은 야추 항목을 연속 클릭","사용 완료 카테고리 검사 → 두 번째 요청 거절",PINK),("정산 중 예외","우승자 저장 중 DB 오류","@Transactional 롤백 → 일부 사용자만 지급되는 상태 방지",YELLOW),("큐 포화","128개를 넘는 변경 요청 유입","AbortPolicy → 빠른 실패와 ‘잠시 후 재시도’ 응답",BLUE)]
for i,(a,b,c,d) in enumerate(scenarios):
    y=1.48+i*1.29; rect(s,.68,y,2.0,.95,PANEL,d); text(s,a,.83,y+.3,1.7,.25,14,d,True,align=PP_ALIGN.CENTER)
    rect(s,2.92,y,3.75,.95,PANEL,LINE); text(s,b,3.18,y+.24,3.25,.42,12,WHITE,True)
    text(s,"→",6.72,y+.25,.4,.35,19,MUTED,True,"Arial",PP_ALIGN.CENTER)
    rect(s,7.2,y,5.45,.95,PANEL2,LINE); text(s,c,7.48,y+.22,4.9,.48,12,MUTED)
text(s,"설계 의도: 성공 경로보다 ‘동시에 요청되거나 중간에 실패했을 때’의 상태를 먼저 정의",1.4,6.68,10.5,.28,14,WHITE,True,align=PP_ALIGN.CENTER)

# 20. Scale roadmap
s=base("현재 구조의 한계와 확장 전략", "DEEP DIVE · SCALABILITY", 19)
left=[("대상","동일 LAN · 최대 6명"),("동기화","1초 폴링"),("상태","단일 서버 메모리"),("명령","단일 Worker Queue")]
right=[("실시간 통신","WebSocket / SSE"),("분산 상태","Redis + 세션 공유"),("동시성","DB Lock / Optimistic Lock"),("처리량","게임 Room별 파티셔닝")]
rect(s,.7,1.5,5.55,4.95,PANEL,LINE); text(s,"현재 선택",1.0,1.83,2.2,.35,19,MINT,True)
rect(s,7.05,1.5,5.55,4.95,PANEL,LINE); text(s,"확장 단계",7.35,1.83,2.2,.35,19,PINK,True)
for i,(a,b) in enumerate(left):
    text(s,a,1.05,2.55+i*.8,1.3,.25,13,MUTED,True); text(s,b,2.4,2.55+i*.8,3.2,.3,14,WHITE,True)
for i,(a,b) in enumerate(right):
    text(s,a,7.4,2.55+i*.8,1.45,.25,13,MUTED,True); text(s,b,8.9,2.55+i*.8,3.1,.3,14,WHITE,True)
text(s,"현재 요구사항에는 단순성과 일관성이 유리하지만, 수평 확장 시 인메모리 상태와 단일 큐가 병목",1.2,6.72,10.95,.28,13,YELLOW,True,align=PP_ALIGN.CENTER)

# 21. AI generation
s=base("AI 문제 출제: 근거를 먼저 고정", "DEEP DIVE · AI GENERATION", 20)
add_image(s,GAME/"LLM WIKI  를기반으로 OPenAI 가문제를 출제 하는 화면.png",.65,1.5,4.05,4.95)
code(s,"List<WikiExcerpt> excerpts = pickExcerpts(3);\n\n아래 LLM Wiki 발췌문만 근거로\n한국어 프로그래밍 퀴즈 한 문제를 만들어라.\n사실이 애매하거나 발췌문으로 검증할 수 없는\n내용은 묻지 마라.\nsource는 제공된 SOURCE 경로를 그대로 사용한다.",4.98,1.5,4.25,3.08,"GROUNDING PROMPT")
rich_lines(s,[("문서 선택  ","_meta·index·log 제외 후 Markdown 3개 무작위 선택",MINT),("길이 제한  ","발췌문 크기를 제한해 입력 비용 통제",PINK),("유형 다양화  ","concept·code_output·debugging·scenario·true_false",YELLOW),("최근 이력  ","최근 20개 핵심 문장과 비교해 반복 감소",BLUE)],9.5,1.62,3.0,3.55,11.5)
text(s,"게임을 만든 이유: 문제 풀이 보상을 공통 크레딧과 연결해 AI 기능이 플랫폼 안에서 지속적으로 소비되도록 설계",4.98,5.15,7.5,1.0,14,WHITE,True)

# 22. AI validation
s=base("AI 응답은 서버 검증을 통과해야 한다", "DEEP DIVE · AI VALIDATION", 21)
code(s,'"type": "object",\n"additionalProperties": false,\n"required": ["category", "type", "prompt",\n             "choices", "answer", "explanation", "source"],\n"choices": { "minItems": 2, "maxItems": 4 },\n"answer":  { "minimum": 0, "maximum": 3 }',.7,1.5,5.45,3.35,"STRUCTURED OUTPUTS · JSON SCHEMA")
code(s,"if (choices.size() < 2 || choices.size() > 4)\n    throw new IllegalArgumentException(...);\nif (answer < 0 || answer >= choices.size())\n    throw new IllegalArgumentException(...);\nif (excerpts.noneMatch(e -> e.source().equals(source)))\n    throw new IllegalArgumentException(\"출처 불일치\");",6.45,1.5,6.2,3.35,"SECOND VALIDATION · SERVER")
rules=[("형식","필수 필드와 선택지 개수",MINT),("범위","정답 인덱스 유효성",PINK),("근거","허용된 위키 경로 일치",YELLOW),("회복","최대 3회 생성 재시도",BLUE)]
for i,(a,b,c) in enumerate(rules):
    x=.8+i*3.08; rect(s,x,5.35,2.75,1.05,PANEL,c); text(s,a,x+.15,5.59,2.45,.22,14,c,True,align=PP_ALIGN.CENTER); text(s,b,x+.15,5.98,2.45,.2,10.5,WHITE,align=PP_ALIGN.CENTER)

# 23. AI storage and transaction
s=base("설명과 출처까지 저장해 서비스 자산으로", "DEEP DIVE · AI PERSISTENCE", 22)
code(s,"@Entity @Table(\n  name = \"quiz_questions\",\n  uniqueConstraints = @UniqueConstraint(columnNames=\"signature\")\n)\nclass QuizQuestionEntity {\n  String prompt; List<String> choices; int answerIndex;\n  String explanation; String source; Instant createdAt;\n}",.7,1.5,5.2,4.75,"QUIZQUESTIONENTITY")
flow2=[("생성","문제·정답·설명·출처",MINT),("정규화","문장 signature 생성",PINK),("중복 확인","findBySignature",YELLOW),("저장","문제 + 선택지 테이블",BLUE),("재사용","AI OFF 시 DB에서 출제",MINT)]
for i,(a,b,c) in enumerate(flow2):
    y=1.48+i*1.02; rect(s,6.25,y,2.05,.76,PANEL,c); text(s,a,6.4,y+.24,1.75,.2,13,c,True,align=PP_ALIGN.CENTER); text(s,"→",8.35,y+.18,.35,.3,16,MUTED,True,"Arial",PP_ALIGN.CENTER); rect(s,8.72,y,3.85,.76,PANEL2,LINE); text(s,b,8.95,y+.21,3.4,.28,12,WHITE,True)
text(s,"정답 제출도 @Transactional: 활성 문제 ID 확인 → 선택지 검증 → 정답 보상 저장",6.25,6.02,6.25,.42,13,YELLOW,True)
text(s,"생성 비용을 일회성으로 바꾸고, 장애·비활성화 상황에서도 퀴즈 서비스를 유지",3.25,6.7,7.4,.28,14,WHITE,True,align=PP_ALIGN.CENTER)

# Planning background — make the engineering intent explicit.
planning_slide = base("기획 배경: 게임보다 상태 변화가 많은 서비스", "01 · Planning background", 4)
text(planning_slide, "출발점", .75, 1.52, 1.2, .28, 12, MINT, True)
text(planning_slide, "게임은 짧은 시간에 ‘참가 → 차감 → 판정 → 정산’이 반복됩니다.", 2.05, 1.45, 9.8, .58, 17, WHITE, True)
rect(planning_slide, .75, 2.32, 11.55, .012, LINE, LINE, False)
planning_cols = [
    ("기획 의도", "반복되는 크레딧 변화와 게임 상태 전이를 통해 트랜잭션, 동시 요청, 상태 동기화를 구현하고자 했습니다.", MINT),
    ("핵심 문제", "여러 사용자의 요청이 같은 시점에 도착해도 중복 차감과 서로 다른 결과가 발생하지 않아야 합니다.", BLUE),
    ("최종 범위", "레이싱·야추·행맨·마피아가 하나의 플레이어와 크레딧을 공유하며 서버가 결과를 최종 판정합니다.", YELLOW),
]
for i, (head, body, c) in enumerate(planning_cols):
    x = .75 + i * 3.9
    text(planning_slide, f"0{i+1}", x, 2.72, .4, .25, 11, c, True, "Arial")
    text(planning_slide, head, x, 3.12, 3.45, .35, 17, WHITE, True)
    text(planning_slide, body, x, 3.72, 3.45, 1.35, 13.5, MUTED)
text(planning_slide, "서비스 경계", .75, 5.65, 1.3, .25, 12, PINK, True)
text(planning_slide, "크레딧은 서비스 내부에서만 사용하는 게임 점수이며, 동일 LAN 최대 6명을 실제 목표 범위로 설정", 2.15, 5.58, 9.8, .38, 14, WHITE, True)
text(planning_slide, "핵심 질문  ·  두 사용자가 동시에 마지막 크레딧으로 참가하면 서버 상태는 어떻게 설명되어야 하는가?", .75, 6.48, 11.55, .42, 15, BLUE, True, align=PP_ALIGN.CENTER)
planning_id = prs.slides._sldIdLst[-1]

# Credit lifecycle and transaction boundary.
economy_slide = base("크레딧 정산의 전체 생명주기", "03 · Credit lifecycle", 7)
stages = [
    ("1", "요청", "player · game · stake", MINT),
    ("2", "검증", "잔액 · 허용 금액 · 중복", BLUE),
    ("3", "차감", "Player.debit + save", YELLOW),
    ("4", "진행", "참가자·턴·pot 상태", PINK),
    ("5", "판정", "서버 결과 단일 생성", BLUE),
    ("6", "정산", "winner.credit + save", MINT),
]
for i, (n, head, body, c) in enumerate(stages):
    x = .55 + i * 2.08
    text(economy_slide, n, x, 1.55, .32, .28, 11, c, True, "Arial")
    text(economy_slide, head, x, 1.92, 1.72, .3, 16, WHITE, True)
    text(economy_slide, body, x, 2.38, 1.72, .56, 11.5, MUTED)
    if i < 5: text(economy_slide, "→", x + 1.72, 2.0, .34, .3, 16, MUTED, True, "Arial", PP_ALIGN.CENTER)
rect(economy_slide, .7, 3.32, 5.85, 2.65, PANEL, LINE)
text(economy_slide, "DB 트랜잭션 안에서 보장", 1.0, 3.63, 4.9, .3, 17, MINT, True)
rich_lines(economy_slide, [
    ("원자성  ", "차감·보상 저장 중 예외가 발생하면 DB 변경 롤백", MINT),
    ("검증  ", "허용 stake와 현재 잔액을 서버 값으로 다시 확인", BLUE),
    ("정산  ", "게임별 참가비와 보상을 동일 Player 크레딧으로 통합", YELLOW),
], 1.0, 4.18, 5.0, 1.45, 12.5)
rect(economy_slide, 6.82, 3.32, 5.8, 2.65, PANEL, LINE)
text(economy_slide, "DB 트랜잭션만으로 부족한 부분", 7.12, 3.63, 5.0, .3, 17, PINK, True)
rich_lines(economy_slide, [
    ("메모리 상태  ", "턴·참가자·pot은 JPA 롤백 대상이 아님", PINK),
    ("경쟁 조건  ", "명령 큐와 synchronized로 변경 순서를 통제", BLUE),
    ("확장 시  ", "Redis/DB 상태화와 낙관적 락·멱등 키가 필요", YELLOW),
], 7.12, 4.18, 5.0, 1.45, 12.5)
text(economy_slide, "한 번의 게임 사건이 어디서 시작하고 어디까지 함께 성공해야 하는지를 트랜잭션 경계로 정의", 1.05, 6.55, 11.2, .35, 14, WHITE, True, align=PP_ALIGN.CENTER)
economy_id = prs.slides._sldIdLst[-1]

# Concrete concurrency sequence.
concurrency_slide = base("동시에 들어온 참가 요청을 어떻게 처리했는가", "Deep dive · Concurrency", 21)
text(concurrency_slide, "상황: A와 B가 거의 동시에 마지막 크레딧으로 참가", .75, 1.48, 8.6, .35, 17, WHITE, True)
actors = [("Client A", .85, MINT), ("Client B", .85, BLUE), ("Command Queue", 4.75, YELLOW), ("Service / DB", 8.72, PINK)]
for idx, (name, x, c) in enumerate(actors):
    y = 2.0 if idx != 1 else 2.92
    rect(concurrency_slide, x, y, 3.05, .65, PANEL, c)
    text(concurrency_slide, name, x + .18, y + .2, 2.7, .24, 13, c, True, "Arial", PP_ALIGN.CENTER)
text(concurrency_slide, "POST join(A)", 1.0, 3.94, 2.7, .25, 12, WHITE, True)
text(concurrency_slide, "POST join(B)", 1.0, 4.62, 2.7, .25, 12, WHITE, True)
text(concurrency_slide, "① enqueue", 3.55, 3.94, 1.1, .25, 11, MUTED)
text(concurrency_slide, "② enqueue", 3.55, 4.62, 1.1, .25, 11, MUTED)
rect(concurrency_slide, 4.75, 3.72, 3.05, 1.35, PANEL2, LINE)
text(concurrency_slide, "A 실행 완료 후 B 실행", 5.0, 4.14, 2.55, .42, 14, WHITE, True, align=PP_ALIGN.CENTER)
text(concurrency_slide, "③ 최신 잔액 검증", 7.85, 3.94, 1.25, .25, 11, MUTED)
text(concurrency_slide, "④ B 재검증", 7.85, 4.62, 1.25, .25, 11, MUTED)
rect(concurrency_slide, 9.15, 3.72, 3.0, 1.35, PANEL, LINE)
text(concurrency_slide, "A: 차감 성공\nB: 잔액 부족 거절", 9.4, 4.0, 2.5, .65, 14, WHITE, True, align=PP_ALIGN.CENTER)
text(concurrency_slide, "순서 보장", .85, 5.72, 1.2, .25, 12, MINT, True)
text(concurrency_slide, "단일 Worker Queue", 2.08, 5.68, 2.2, .3, 13, WHITE, True)
text(concurrency_slide, "임계 구역", 4.48, 5.72, 1.2, .25, 12, BLUE, True)
text(concurrency_slide, "synchronized service", 5.72, 5.68, 2.25, .3, 13, WHITE, True)
text(concurrency_slide, "DB 원자성", 8.25, 5.72, 1.2, .25, 12, YELLOW, True)
text(concurrency_slide, "@Transactional", 9.48, 5.68, 2.2, .3, 13, WHITE, True)
text(concurrency_slide, "현재 선택은 처리량보다 설명 가능한 상태를 우선하며, 확장 시 Room별 큐로 병렬성을 확보", 1.05, 6.58, 11.2, .34, 14, PINK, True, align=PP_ALIGN.CENTER)
concurrency_id = prs.slides._sldIdLst[-1]

# Promotion UI evidence using the newly supplied photo.
ad_slide = base("외부 프로모션을 서비스 UI에 연결", "Feature · Promotion", 15)
add_image(ad_slide, GAME / "배너광고 이미지.png", .7, 1.48, 5.7, 1.18)
add_image(ad_slide, ROOT / "모달광고.png", .7, 2.88, 5.7, 3.98)
text(ad_slide, "배너 + 세션 모달 광고", 6.82, 1.55, 5.25, .4, 20, WHITE, True)
text(ad_slide, "실제 브라우저에서 캡처한 배너와 모달 화면을 원본 비율로 배치했습니다.", 6.82, 2.08, 5.25, .58, 13.5, MUTED)
rich_lines(ad_slide, [
    ("소재  ", "로컬 image1·image2를 사용해 외부 이미지 만료 방지", MINT),
    ("노출  ", "배너는 상시, 모달은 sessionStorage 기준 세션당 1회", BLUE),
    ("반응형  ", "원본 1:1 비율 유지, 화면 높이에 맞춰 최대 크기 제한", YELLOW),
    ("이동  ", "원본 Instagram 게시물 링크를 새 탭으로 연결", PINK),
    ("접근성  ", "대체 텍스트·닫기 버튼·dialog 속성 적용", MINT),
], 6.82, 2.95, 5.15, 2.75, 12.5)
text(ad_slide, "게임 기능 외에도 실제 서비스에서 필요한 콘텐츠 노출과 외부 캠페인 연결 흐름을 구현", 6.82, 6.18, 5.15, .58, 14, WHITE, True)
ad_id = prs.slides._sldIdLst[-1]

# Measured concurrency result — evidence for the design claims.
load_test_slide = base("최대 6명 동시 요청 검증 결과", "Validation · Measured result", 23)
add_image(load_test_slide, ROOT / "부하테스트 이미지.png", .7, 1.48, 7.25, 1.88)
rect(load_test_slide, 8.28, 1.48, 4.32, 1.88, PANEL, LINE)
text(load_test_slide, "FINAL RESULT", 8.58, 1.78, 3.7, .22, 10, MUTED, True, "Arial")
text(load_test_slide, "PASS", 8.58, 2.15, 3.7, .48, 28, MINT, True, "Arial")
text(load_test_slide, "동일 LAN 최대 6명 요구사항 기준", 8.58, 2.74, 3.5, .25, 11.5, MUTED)
results = [
    ("6 / 6", "동시 입장 성공", MINT),
    ("6 / 6", "야추 동시 참가", BLUE),
    ("5 / 5", "중복 요청 거절", YELLOW),
    ("60 / 60", "상태 조회 성공", PINK),
]
for i, (value, label, c) in enumerate(results):
    x = .72 + i * 3.0
    rect(load_test_slide, x, 3.72, 2.72, 1.18, PANEL, LINE)
    text(load_test_slide, value, x + .16, 3.94, 2.4, .36, 22, c, True, "Arial", PP_ALIGN.CENTER)
    text(load_test_slide, label, x + .16, 4.48, 2.4, .22, 11, MUTED, True, align=PP_ALIGN.CENTER)
text(load_test_slide, "응답 시간", .75, 5.38, 1.2, .25, 13, WHITE, True)
latencies = [("평균", "30.37 ms"), ("p95", "79.85 ms"), ("최대", "107.17 ms"), ("실패", "0건")]
for i, (label, value) in enumerate(latencies):
    x = 2.1 + i * 2.45
    text(load_test_slide, label, x, 5.38, .72, .24, 11, MUTED, True)
    text(load_test_slide, value, x + .72, 5.34, 1.55, .3, 14, WHITE, True, "Arial")
text(load_test_slide, "검증 조건", .75, 6.08, 1.2, .25, 12, BLUE, True)
text(load_test_slide, "6명 동시 참가 시 전원 100 크레딧 1회만 차감 · 동일 사용자 중복 참가 5건 전부 거절", 2.1, 6.02, 10.15, .38, 13.5, WHITE, True)
text(load_test_slide, "※ 대규모 부하 성능이 아니라 실제 목표 범위의 동시성·일관성 검증", .75, 6.66, 11.5, .25, 11, MUTED, False, align=PP_ALIGN.CENTER)
load_test_id = prs.slides._sldIdLst[-1]

# Test harness implementation — show how the measured result was produced.
test_code_slide = base("Python으로 반복 가능한 동시 요청 테스트 구성", "Validation · Test harness", 24)
code(test_code_slide, '''def parallel(jobs, workers=12):
    with ThreadPoolExecutor(max_workers=workers) as pool:
        futures = [pool.submit(call, *job) for job in jobs]
        return [future.result() for future in futures]

# 6명의 참가 요청을 동시에 전송
yacht_jobs = [
    (base_url, "/api/yacht/join", token, {"bet": 100})
    for token in tokens
]
yacht_responses = parallel(yacht_jobs, workers=6)

join_success = sum(r.ok for r in yacht_responses)
credits_once = all(
    state.data["me"]["credits"]
      == initial_credits[state.data["me"]["id"]] - 100
    for state in after_join_states
)''', .7, 1.48, 6.05, 4.95, "CONCURRENT REQUEST · VERIFICATION")
steps = [
    ("01", "동시 시작", "ThreadPoolExecutor로 6개의 HTTP 요청을 병렬 전송", MINT),
    ("02", "결과 확인", "응답 성공 수, 야추 로비 인원과 최종 크레딧 비교", BLUE),
    ("03", "중복 검증", "같은 토큰의 참가 요청 5건을 보내 모두 거절되는지 확인", YELLOW),
    ("04", "조회 측정", "상태 API 60건의 성공률과 평균·p95·최대 시간 집계", PINK),
    ("05", "보고서 저장", "실행 조건과 결과를 Markdown으로 자동 기록", MINT),
]
for i, (n, head, body, c) in enumerate(steps):
    y = 1.5 + i * 1.02
    text(test_code_slide, n, 7.12, y + .14, .42, .22, 11, c, True, "Arial")
    text(test_code_slide, head, 7.72, y + .1, 1.45, .28, 14, WHITE, True)
    text(test_code_slide, body, 9.25, y + .08, 3.15, .48, 11.5, MUTED)
    rect(test_code_slide, 7.12, y + .73, 5.25, .012, LINE, LINE, False)
text(test_code_slide, "실행", .72, 6.67, .6, .22, 11, BLUE, True)
text(test_code_slide, "python scripts/local_arcade_concurrency_test.py --join-code <6자리 코드>", 1.4, 6.62, 7.1, .28, 11.5, WHITE, True, "Consolas")
text(test_code_slide, "산출물  reports/local-arcade-concurrency-*.md", 8.65, 6.62, 3.7, .28, 10.5, MUTED, False, "Consolas", PP_ALIGN.RIGHT)
test_code_id = prs.slides._sldIdLst[-1]

# Interview follow-up slides. Each is inserted immediately after the related
# implementation page at the end of the build.
qa_specs = [
    ("아키텍처", "System Architecture", [
        ("왜 클라이언트 상태를 최종 값으로 사용하지 않았나요?", "멀티플레이에서는 각 브라우저의 시점과 입력이 다를 수 있습니다. 클라이언트는 입력과 표현만 담당하고, 턴·결과·크레딧은 서버 값으로 다시 검증해 하나의 기준을 유지했습니다."),
        ("단일 서버 메모리에 게임 상태를 둔 이유와 한계는 무엇인가요?", "최대 6명 LAN 환경에서는 구현 복잡도와 응답 비용이 낮다는 장점이 있습니다. 대신 재시작과 수평 확장에 취약하므로 규모가 커지면 Room 상태와 세션을 Redis 또는 DB로 외부화해야 합니다."),
        ("서비스 계층이 한 서버에 모여 있어 결합도가 높지 않나요?", "현재는 배포 단위가 하나인 모듈형 모놀리스에 가깝습니다. 게임별 Service는 분리했지만 Player와 Credit은 공유합니다. 독립 확장이 필요해지는 시점에 이벤트와 명확한 경계를 기준으로 분리할 계획입니다."),
    ]),
    ("크레딧 정산", "크레딧 정산의 전체 생명주기", [
        ("@Transactional을 붙이면 게임 상태도 모두 롤백되나요?", "아닙니다. JPA가 관리하는 DB 변경만 롤백되고 Java 메모리의 참가자·턴·pot은 대상이 아닙니다. 이 한계를 슬라이드에 명시했고, 확장 시 게임 상태도 영속 저장소로 이동해야 합니다."),
        ("명령 큐와 synchronized를 함께 쓰는 이유는 무엇인가요?", "큐는 변경 명령의 순서와 적재량을 통제하고, synchronized는 서비스 임계 구역을 보호합니다. 현재는 방어적으로 함께 사용하지만 Room별 큐가 서비스 호출 경로를 완전히 통제하면 임계 범위를 줄일 수 있습니다."),
        ("차감과 보상의 트랜잭션 경계가 너무 길지 않나요?", "각 API 호출 안의 DB 변경만 짧게 묶었습니다. 게임 전체 시간을 하나의 트랜잭션으로 유지하지 않습니다. 참가 차감과 결과 정산은 서로 다른 사건으로 분리되어 있습니다."),
    ]),
    ("게임 상태와 결과 동기화", "모든 플레이어가 같은 결과를 공유하는 레이싱", [
        ("야추에서 화면의 주사위 값을 요청으로 보내면 조작할 수 있지 않나요?", "요청 값은 판정 자료로 신뢰하지 않고 서버가 보유한 배열과 일치하는지만 확인합니다. 점수 계산도 서버 주사위로 수행하므로 수정된 클라이언트 값으로 점수를 만들 수 없습니다."),
        ("레이싱 결과의 공정성은 어떻게 보장하나요?", "서버의 SecureRandom으로 결과를 한 번 생성해 모든 사용자에게 동일하게 제공합니다. 현재는 내부 게임 점수 범위이며 감사 가능한 공정성이 요구되면 결과 로그와 seed 공개 또는 commit-reveal 방식이 추가로 필요합니다."),
        ("클라이언트가 같은 결과 애니메이션을 여러 번 실행할 수 있지 않나요?", "각 결과에 raceId를 부여하고 브라우저가 마지막으로 본 ID와 비교합니다. 새로운 ID일 때만 애니메이션을 실행해 폴링 중복 표시를 막았습니다."),
    ]),
    ("AI 퀴즈 파이프라인", "근거 기반 AI 퀴즈 파이프라인", [
        ("JSON Schema를 사용하면 AI의 사실 오류도 없어지나요?", "아닙니다. Schema는 필드와 자료형만 보장합니다. 별도로 제공한 위키 경로와 source 일치 여부, 정답 인덱스, 선택지 범위를 서버에서 검사하며 사실성은 근거 제한과 향후 관리자 검수로 보완합니다."),
        ("로컬 위키 문서에 잘못된 내용이 있으면 어떻게 되나요?", "현재 구조는 위키를 신뢰 가능한 내부 자료로 가정합니다. 따라서 문서 품질이 상한선이며, 운영 환경에서는 문서 승인 상태·버전·출처 해시를 저장해 어떤 근거로 생성됐는지 추적해야 합니다."),
        ("OpenAI API 장애나 비용 문제는 어떻게 처리하나요?", "생성 결과를 DB에 저장해 재사용하고, 관리자가 생성을 끄면 저장 문제를 출제합니다. 따라서 외부 API가 없어도 기존 퀴즈 흐름은 유지됩니다."),
    ]),
    ("프로모션 UI", "외부 프로모션을 서비스 UI에 연결", [
        ("광고 기능이 프로젝트 핵심 기술과 어떤 관련이 있나요?", "핵심 동시성 기능은 아니지만 실제 서비스에서 필요한 외부 콘텐츠 노출, 반응형 이미지, 접근성, 링크 연결과 노출 정책을 적용한 운영 기능입니다. 기술 핵심과 구분해 한 장으로만 설명합니다."),
        ("Instagram 이미지를 직접 요청하지 않고 로컬에 둔 이유는 무엇인가요?", "외부 CDN 주소는 만료되거나 접근 정책에 따라 실패할 수 있습니다. 허용된 캠페인 이미지를 로컬 정적 자산으로 관리해 화면 렌더링을 안정화했습니다."),
        ("sessionStorage로 노출을 제한하면 사용자 단위 제어가 가능한가요?", "브라우저 탭 세션 단위 제어만 가능합니다. 현재 요구에는 충분하지만 계정별 빈도 제한이나 노출 통계가 필요하면 서버 이벤트와 사용자별 저장이 필요합니다."),
    ]),
    ("통신 방식", "통신 구조: 현재 구현과 확장 경계", [
        ("실시간 게임인데 왜 WebSocket을 사용하지 않았나요?", "최대 6명, 1초 수준 동기화라는 요구에서는 HTTP 폴링이 구현과 재연결 처리가 단순했습니다. 즉각적인 push가 필요한 규모가 되면 WebSocket과 Room 구독 구조로 전환할 계획입니다."),
        ("기본 약 12 req/s는 실제 측정값인가요?", "6명 × 2개 API × 초당 1회로 계산한 요청 모델입니다. 실제 자동 테스트에서는 상태 조회 60건의 성공률과 응답 시간을 별도로 측정했으며 두 수치를 구분해 표기했습니다."),
        ("최대 1초의 오래된 상태로 명령을 보내면 문제가 없나요?", "클라이언트 화면은 최대 1초 늦을 수 있습니다. 그래서 변경 요청 시 현재 턴·참가 여부·서버 주사위와 잔액을 다시 검사하고, 오래된 요청은 거절한 뒤 최신 상태를 다시 조회합니다."),
    ]),
    ("자동 동시 요청 테스트", "Python으로 반복 가능한 동시 요청 테스트 구성", [
        ("Python ThreadPoolExecutor가 정말 동시에 요청을 보내나요?", "HTTP 요청은 I/O 작업이라 여러 Worker가 겹쳐 수행됩니다. 다만 완전히 같은 CPU 시각을 보장하는 시험은 아니며, 현재 테스트는 서버에 중첩 요청을 만들어 기능 일관성을 검증하는 목적입니다."),
        ("60건과 6명으로 성능을 주장할 수 있나요?", "대규모 성능을 주장할 수 없습니다. 최대 6명이라는 실제 요구에서 성공률과 중복 차감 방지를 확인한 기능·동시성 검증입니다. p95도 이 실행 표본의 관찰값으로만 사용합니다."),
        ("테스트가 무엇을 자동으로 판정하나요?", "6명 입장, 로비 인원, 사용자별 100 크레딧 1회 차감, 중복 요청 5건 거절, 상태 조회 실패 0건을 코드로 비교하고 Markdown 보고서를 생성합니다."),
    ]),
    ("실패와 트랜잭션", "실패 시나리오로 보는 데이터 일관성", [
        ("정산 중 DB 예외가 나면 일부 사용자만 보상받을 수 있지 않나요?", "동일 정산 메서드의 JPA 변경은 하나의 트랜잭션으로 묶여 예외 시 롤백됩니다. 다만 이미 바뀐 메모리 게임 상태는 자동 복구되지 않으므로 상태 외부화나 보상 트랜잭션이 다음 과제입니다."),
        ("중복 클릭은 프론트 버튼을 비활성화하면 충분하지 않나요?", "브라우저 조작, 네트워크 재시도와 직접 API 호출은 막을 수 없습니다. 사용 완료 카테고리와 참가 여부를 서버에서 다시 검사해야 최종 일관성을 보장할 수 있습니다."),
        ("큐가 가득 차면 사용자 경험은 어떻게 되나요?", "128개를 넘으면 무한 대기 대신 즉시 실패하도록 설계했습니다. 현재는 재시도 안내를 반환하고, 운영 단계에서는 429 응답·Retry-After·지표 수집을 추가하는 것이 적절합니다."),
    ]),
    ("확장 전략", "현재 구조의 한계와 확장 전략", [
        ("서버를 두 대로 늘리면 현재 구조가 바로 동작하나요?", "동작하지 않습니다. 세션과 게임 상태가 각 JVM에 나뉘고 단일 큐도 서버별로 분리됩니다. Redis 세션·Room 상태·Pub/Sub와 일관된 라우팅 또는 분산 락이 필요합니다."),
        ("Room별 큐로 나누면 같은 플레이어 요청의 순서는 어떻게 보장하나요?", "roomId를 파티션 키로 사용해 같은 Room 명령은 동일 소비자에서 순서대로 처리하고 서로 다른 Room만 병렬 처리합니다. 메시지 ID와 멱등 키로 중복 실행도 방지해야 합니다."),
        ("낙관적 락만 적용하면 동시성 문제가 해결되나요?", "DB 행 충돌은 감지할 수 있지만 인메모리 턴 상태와 메시지 순서는 해결하지 못합니다. 상태 저장 위치, 명령 순서, 재시도 정책을 함께 설계해야 합니다."),
    ]),
    ("AI 저장과 검증", "설명과 출처까지 저장해 서비스 자산으로", [
        ("signature unique 제약으로 의미가 같은 문제도 막을 수 있나요?", "현재 signature는 정규화된 문장 중복을 막는 수준이라 표현이 다른 의미 중복은 남을 수 있습니다. 임베딩 유사도나 핵심 개념 키를 추가하면 의미 중복을 더 줄일 수 있습니다."),
        ("source 경로가 일치하면 설명이 정확하다고 볼 수 있나요?", "볼 수 없습니다. 경로 검사는 허용되지 않은 출처 사용을 막을 뿐입니다. 문장 인용 범위 저장, 근거 문장 대조와 관리자 승인 상태가 추가되어야 품질을 높일 수 있습니다."),
        ("문제 저장과 정답 보상은 어떻게 연결되나요?", "활성 문제 ID와 선택지 범위를 서버가 확인한 뒤 정답이면 같은 트랜잭션에서 크레딧을 저장합니다. 클라이언트가 정답 여부나 보상 금액을 결정하지 않습니다."),
    ]),
]
qa_insertions = []
for i, (topic, target_title, questions) in enumerate(qa_specs):
    qa = interview_slide(topic, questions, 40 + i)
    qa_insertions.append((prs.slides._sldIdLst[-1], target_title))

detail_specs = [
    ("요청 한 건이 처리되는 전체 경로", "System Architecture",
     "조회와 변경 요청이 같은 인증 절차를 거치지만 변경 명령만 큐에 진입합니다.",
     '''GET /api/state
  -> auth(token)
  -> service.state(player)

POST /api/yacht/score
  -> auth(token)
  -> queue.run(command)
  -> validate current turn
  -> @Transactional save
  -> response state''',
     [("인증", "Bearer 토큰을 Player ID로 변환하고 lastSeen을 갱신", MINT),
      ("경로 분리", "조회는 즉시 처리하고 변경만 순서 제어", BLUE),
      ("도메인 검증", "Controller가 아닌 Service가 현재 상태를 확인", YELLOW),
      ("응답 동기화", "명령 결과와 함께 최신 서버 상태를 반환", PINK)],
     "GameController의 GET/POST 매핑과 GameCommandQueue.run 호출 경로"),
    ("공유 상태의 소유권을 어디에 두었는가", "System Architecture",
     "상태마다 수명과 일관성 요구가 달라 저장 위치를 구분했습니다.",
     '''Browser
  UI tab / animation / input

Server memory
  turn / lobby / chat / race result

MySQL
  player / credits / quiz question

External
  wiki source / AI response''',
     [("브라우저", "잃어도 복구 가능한 표현 상태만 보유", MINT),
      ("서버 메모리", "짧은 게임 세션의 빠른 공유 상태", BLUE),
      ("DB", "재시작 후에도 필요한 크레딧과 문제", YELLOW),
      ("외부 입력", "위키와 AI 응답은 검증 후 내부 데이터로 변환", PINK)],
     "서버 재조회만으로 화면을 복구할 수 있도록 최종 판정 상태를 클라이언트 밖에 배치"),
    ("크레딧 변경에서 지켜야 할 불변 조건", "크레딧 정산의 전체 생명주기",
     "정산 코드는 기능보다 먼저 깨지면 안 되는 조건을 정의했습니다.",
     '''void debit(long amount) {
  if (amount < 10 || amount > 5000)
    throw new IllegalArgumentException(...);
  if (credits < amount)
    throw new IllegalArgumentException(...);
  credits -= amount;
}

// invariant
credits >= 0
one request == one debit''',
     [("음수 방지", "현재 잔액보다 큰 차감 요청을 서버에서 거절", MINT),
      ("허용 범위", "게임별 선택 가능한 금액을 다시 검사", BLUE),
      ("단일 차감", "중복 참가 여부를 확인한 뒤 한 번만 저장", YELLOW),
      ("정산 합계", "게임 결과의 지급 총액이 규칙과 일치해야 함", PINK)],
     "Player.debit, allowed, joinYacht의 검증 순서와 6명 동시 요청 테스트"),
    ("세션 생성과 비활성 사용자 정리", "User Flow",
     "참여 코드 입장부터 비활성 세션 제거까지 사용자 수명을 서버가 관리합니다.",
     '''String token = UUID.randomUUID().toString();
sessions.put(token, player.getId());

@Scheduled(fixedDelay = 30000)
void cleanupInactivePlayers() {
  var stale = repo.findByLastSeenBefore(
      now.minusSeconds(30));
  sessions.removeIf(stale);
  repo.deleteAll(stale);
}''',
     [("입장 제한", "닉네임 중복과 최대 6명을 서버에서 검사", MINT),
      ("토큰", "브라우저에는 임의 세션 토큰만 전달", BLUE),
      ("활동 갱신", "인증 요청마다 lastSeen을 갱신", YELLOW),
      ("정리", "30초 주기로 세션·게임 참가 상태를 제거", PINK)],
     "GameService.join, auth, cleanupInactivePlayers의 실제 사용자 수명주기"),
    ("야추를 상태 머신으로 해석하기", "턴과 상태를 서버가 관리하는 멀티플레이 야추",
     "버튼별 코드가 아니라 허용 가능한 상태 전이로 게임 규칙을 구성했습니다.",
     '''LOBBY
  join -> ready -> START

TURN
  roll(1..3) -> hold -> score

score
  -> validate category
  -> save points
  -> next turn | FINISHED

FINISHED
  settle -> restart | leave''',
     [("Lobby", "참가 금액 통일과 최소 인원을 검증", MINT),
      ("Turn", "현재 플레이어와 굴림 횟수를 서버가 보유", BLUE),
      ("Score", "사용 완료 항목과 서버 주사위를 검증", YELLOW),
      ("Finish", "공동 우승·기권·환불까지 종료 상태로 처리", PINK)],
     "SocialGameService의 joinYacht, roll, score, cancelYacht, leave 메서드"),
    ("레이싱 결과가 모든 화면에 동일하게 도착하는 과정", "모든 플레이어가 같은 결과를 공유하는 레이싱",
     "랜덤 생성과 애니메이션을 분리해 판정은 한 번, 표현은 각 브라우저에서 수행합니다.",
     '''POST /race/start
  Server: shuffle order once
  Server: save RaceResult(raceId, order)
  Client A: poll -> new raceId -> animate
  Client B: poll -> new raceId -> animate
  Client C: poll -> same order

if (raceId != seenRaceId) {
  seenRaceId = raceId;
  playAnimation(order);
}''',
     [("단일 판정", "서버가 순위와 보상을 한 번 생성", MINT),
      ("결과 식별", "raceId가 새 결과와 이전 결과를 구분", BLUE),
      ("지연 허용", "도착 시점은 달라도 order는 동일", YELLOW),
      ("표현 분리", "주행 시간은 결과 판정에 영향을 주지 않음", PINK)],
     "GameService.race의 RaceResult와 main.tsx의 seenRace ref 비교"),
    ("마피아 단계와 채팅 권한 처리", "게임 규칙을 상태 머신으로 확장",
     "같은 채팅 기능도 현재 단계와 역할에 따라 접근 규칙이 달라집니다.",
     '''String scope = phase.equals("NIGHT")
    ? "MAFIA" : "DAY";

if (scope.equals("MAFIA")
    && !role.equals("MAFIA"))
  throw new IllegalArgumentException(...);

chat.add(new Message(
  nickname, text, scope, round));''',
     [("단계", "DAY·NIGHT·FINISHED를 서버 시간으로 전환", MINT),
      ("역할", "밤에는 허용된 역할만 행동과 채팅 가능", BLUE),
      ("생존", "탈락 사용자의 투표·채팅 요청을 거절", YELLOW),
      ("범위", "메시지에 phase와 round를 저장해 노출을 분리", PINK)],
     "MafiaGameService.chat, act, advance, requireAlive의 서버 권한 검사"),
    ("시간 기반 캐릭터 이동 계산", "시간 기반 2D 캐릭터 컨트롤",
     "프레임 수가 아니라 경과 시간으로 이동량을 계산해 환경 차이를 줄였습니다.",
     '''const dt = Math.min(
  (now - previous) / 1000,
  0.032
);

x += direction * velocity * dt;
vy -= GRAVITY * dt;
y += vy * dt;

if (landing) {
  y = platformTop; vy = 0;
}''',
     [("Delta time", "초당 속도를 실제 경과 시간에 곱함", MINT),
      ("상한", "긴 프레임 이후 이동 폭증을 32ms로 제한", BLUE),
      ("충돌", "이전·현재 위치로 발판 통과 여부를 확인", YELLOW),
      ("애니메이션", "입력과 grounded 상태로 idle·run·jump 전환", PINK)],
     "requestAnimationFrame tick과 dt 기반 위치·중력·착지 계산"),
    ("공용 채팅과 관리자 변경의 검증 경계", "공용 채팅과 크레딧 운영 시스템",
     "화면 입력 제한과 별개로 서버가 문자열·권한·변경 단위를 다시 검사합니다.",
     '''String clean = text.strip();
if (clean.isBlank() || clean.length() > 200)
  throw new IllegalArgumentException(...);

chat.addLast(message);
while (chat.size() > 80)
  chat.removeFirst();

admin(code);
allowedAdjustment(amount);''',
     [("입력", "공백·길이 초과 메시지를 서버에서 거절", MINT),
      ("메모리 제한", "최근 80개만 유지해 무한 증가 방지", BLUE),
      ("관리자 인증", "별도 코드와 Header로 운영 API 보호", YELLOW),
      ("변경 단위", "허용된 크레딧 증감 값만 처리", PINK)],
     "SocialGameService.chat과 GameService.admin/grantCredits의 검증 코드"),
    ("배너와 모달의 노출 생명주기", "외부 프로모션을 서비스 UI에 연결",
     "광고 이미지를 표시하는 것뿐 아니라 노출 조건·닫기·외부 이동을 함께 구현했습니다.",
     '''const [promoOpen, setPromoOpen] = useState(
  () => sessionStorage.getItem(
    "arcade-promo-closed") !== "1"
);

function closePromo() {
  setPromoOpen(false);
  sessionStorage.setItem(
    "arcade-promo-closed", "1");
}

target="_blank" rel="noreferrer"''',
     [("초기 노출", "세션 내 닫기 기록이 없을 때만 모달 표시", MINT),
      ("닫기", "상태와 sessionStorage를 함께 갱신", BLUE),
      ("반응형", "원본 비율과 화면 최대 높이를 CSS로 제한", YELLOW),
      ("외부 이동", "새 탭 보안 속성과 대체 텍스트 적용", PINK)],
     "main.tsx의 promoOpen 상태와 styles.css의 배너·모달 미디어 쿼리"),
    ("서버 중심 상태 관리의 방어 계층", "서버 중심 상태 관리",
     "한 가지 기술에 의존하지 않고 서로 다른 실패를 막는 계층을 조합했습니다.",
     '''Request validation
  -> command queue
  -> synchronized service
  -> domain invariant
  -> @Transactional DB write
  -> response from server state''',
     [("Validation", "잘못된 입력과 오래된 상태 요청 거절", MINT),
      ("Ordering", "변경 명령을 하나의 실행 순서로 정렬", BLUE),
      ("Critical section", "공유 Map과 게임 객체의 경쟁 접근 차단", YELLOW),
      ("Persistence", "DB 변경의 원자성과 롤백 보장", PINK)],
     "각 계층은 입력 오류·순서 충돌·메모리 경쟁·DB 부분 저장이라는 서로 다른 문제를 담당"),
    ("현재 테스트 범위와 남은 공백", "검증 가능한 게임 규칙",
     "검증된 영역과 아직 자동화하지 않은 영역을 구분해 다음 작업의 우선순위를 정했습니다.",
     '''Current
  Yacht scoring rules
  Mafia phase / winner rules
  6-user concurrent API script

Next
  Controller integration test
  transaction rollback injection
  frontend component test
  reconnect / stale-state test''',
     [("규칙 테스트", "입력 조합에 따른 결정적 계산 검증", MINT),
      ("API 테스트", "실제 HTTP 동시 요청과 최종 상태 비교", BLUE),
      ("부족한 부분", "예외 주입 기반 롤백 테스트는 아직 미구현", YELLOW),
      ("우선순위", "크레딧과 상태가 함께 바뀌는 경로부터 확대", PINK)],
     "JUnit 도메인 테스트와 reports의 6명 동시 요청 실행 결과를 별도 근거로 관리"),
    ("폴링 요청량을 계산한 방법", "반복 요청이 만드는 서버 부하",
     "요청 모델과 실제 측정값을 분리해 숫자의 의미를 명확히 했습니다.",
     '''Per user / second
  GET /api/state        = 1
  GET /api/social/state = 1

6 users * 2 endpoints
  = 12 read requests / second

Optional
  mafia state = +6 req/s
  admin APIs  = +1 req/s''',
     [("기본 모델", "6명 모두 게임·소셜 상태를 초당 조회", MINT),
      ("조건부 요청", "마피아 참여나 관리자 화면에서 추가", BLUE),
      ("변경 명령", "사용자 행동 시에만 발생해 별도 계산", YELLOW),
      ("측정 구분", "12 req/s는 계산, 60건 결과는 실행 측정", PINK)],
     "frontend setInterval(1000/2000)과 실제 호출 API 목록을 기준으로 산정"),
    ("1초 폴링에서 오래된 상태를 다루는 방법", "통신 구조: 현재 구현과 확장 경계",
     "화면 지연을 허용하되 변경 요청의 정확성은 서버 재검증으로 보호합니다.",
     '''t=0.0  server turn = A
t=0.2  client B last view = A
t=0.7  server advances to B
t=0.8  client A sends stale command
       -> activeTurn(A) rejects
t=1.0  next poll receives turn = B''',
     [("허용 지연", "표현 상태는 최대 약 1초 늦을 수 있음", MINT),
      ("명령 검증", "현재 턴·상태 버전을 서버에서 확인", BLUE),
      ("실패 응답", "오래된 명령을 적용하지 않고 오류 반환", YELLOW),
      ("복구", "다음 조회 또는 명령 응답으로 최신 상태 표시", PINK)],
     "SocialGameService.activeTurn과 클라이언트 1초 refresh 주기의 결합"),
    ("제한된 명령 큐의 내부 동작", "명령 큐로 순서와 과부하를 통제",
     "단일 실행 순서와 제한된 대기열로 일관성과 메모리 상한을 선택했습니다.",
     '''new ThreadPoolExecutor(
  1, 1,
  0L, MILLISECONDS,
  new ArrayBlockingQueue<>(128),
  threadFactory,
  new AbortPolicy()
);

CompletableFuture
  .supplyAsync(command, executor)
  .join();''',
     [("Worker 1", "변경 명령이 완료된 순서를 명확히 함", MINT),
      ("Queue 128", "대기 명령의 메모리 사용량에 상한 설정", BLUE),
      ("Abort", "포화 시 무한 대기보다 빠른 실패 선택", YELLOW),
      ("Exception", "CompletionException에서 원래 도메인 오류 복원", PINK)],
     "GameCommandQueue.java의 ThreadPoolExecutor 생성자와 run 메서드"),
    ("자동 테스트가 결과를 판정하는 방식", "Python으로 반복 가능한 동시 요청 테스트 구성",
     "콘솔 출력만 남기지 않고 전후 상태를 비교해 PASS 조건을 코드로 고정했습니다.",
     '''concurrent_join_passed = (
  join_success == users
  and lobby_count == users
  and credits_once
)

duplicate_passed = (
  duplicate_success == 0
  and duplicate_rejected == 5
  and duplicate_credit_unchanged
)

passed = concurrent_join_passed \
  and duplicate_passed \
  and polling_passed''',
     [("전 상태", "각 Player ID와 초기 크레딧 저장", MINT),
      ("후 상태", "참가 뒤 로비 인원과 크레딧 다시 조회", BLUE),
      ("중복", "동일 토큰 요청 5건의 오류 메시지 확인", YELLOW),
      ("보고서", "조건·성공률·latency를 Markdown으로 저장", PINK)],
     "scripts/local_arcade_concurrency_test.py와 생성된 reports Markdown"),
    ("트랜잭션 경계를 API 사건과 맞춘 이유", "트랜잭션 경계를 게임 결과와 맞추다",
     "긴 게임 전체가 아니라 하나의 참가·점수 확정·정산 사건만 원자적으로 처리합니다.",
     '''@Transactional
Map<String,Object> score(...) {
  YachtGame g = activeTurn(player);
  validateDiceAndCategory(g, request);
  sheet.put(category, calculate(...));

  if (finished)
    finishYacht(g); // credit + save
  else
    advanceTurn(g);

  return state(player);
}''',
     [("시작", "현재 턴과 요청 데이터를 검증한 뒤 변경", MINT),
      ("DB 범위", "해당 호출에서 발생한 크레딧 저장을 함께 처리", BLUE),
      ("종료", "응답 직전에 최신 상태를 다시 구성", YELLOW),
      ("주의", "Map과 YachtGame 변경은 DB rollback 밖에 존재", PINK)],
     "SocialGameService.score의 메서드 트랜잭션과 synchronized 임계 구역"),
    ("실패 이후 상태를 복구하는 기준", "실패 시나리오로 보는 데이터 일관성",
     "실패 종류마다 거절·롤백·재조회·재시도 중 다른 복구 방식을 선택합니다.",
     '''Invalid request
  -> reject / no mutation

DB exception
  -> rollback DB transaction

Stale client state
  -> reject / return latest state

Queue saturation
  -> fail fast / retry later

Inactive player
  -> scheduled cleanup''',
     [("사전 실패", "검증 단계에서 변경 전 거절", MINT),
      ("저장 실패", "트랜잭션 롤백 후 오류 응답", BLUE),
      ("상태 불일치", "서버 값을 기준으로 화면 재동기화", YELLOW),
      ("운영 실패", "포화·비활성 상태를 제한과 정리로 해소", PINK)],
     "중복 참가·점수 중복·DB 예외·큐 포화·비활성 사용자 시나리오를 분리"),
    ("수평 확장을 단계별로 진행하는 이유", "현재 구조의 한계와 확장 전략",
     "기능을 한 번에 분산시키지 않고 실제 병목이 나타나는 순서로 이동합니다.",
     '''Phase 1
  HTTP polling -> WebSocket

Phase 2
  session / room state -> Redis

Phase 3
  one queue -> room partitions

Phase 4
  DB version / idempotency key

Phase 5
  metrics / tracing / autoscale''',
     [("통신", "변경 이벤트 push로 불필요한 조회 감소", MINT),
      ("상태", "서버 간 세션과 Room 상태 공유", BLUE),
      ("처리", "Room별 순서는 유지하며 서로 다른 Room 병렬화", YELLOW),
      ("운영", "충돌·중복·지연을 지표로 관찰", PINK)],
     "현재 병목인 폴링·JVM 상태·단일 큐를 순서대로 외부화하는 로드맵"),
    ("AI 프롬프트에 근거를 고정하는 방법", "AI 문제 출제: 근거를 먼저 고정",
     "모델이 자유롭게 지식을 추가하지 않도록 입력 문서와 출력 역할을 제한했습니다.",
     '''excerpts = pickExcerpts(3);

SYSTEM:
  제공된 발췌문만 근거로 사용한다.
  검증할 수 없는 내용은 묻지 않는다.

INPUT:
  [SOURCE path]
  excerpt text...

OUTPUT:
  prompt / choices / answer
  explanation / source''',
     [("문서 선택", "출제 가능한 Markdown만 후보로 사용", MINT),
      ("입력 제한", "발췌 크기와 문서 개수를 제한", BLUE),
      ("지시", "근거 밖 내용을 질문하지 않도록 명시", YELLOW),
      ("추적", "source를 출력 필수 필드로 요구", PINK)],
     "QuizGameService의 pickExcerpts와 Responses API 입력 프롬프트"),
    ("AI 응답을 두 단계로 검증하는 이유", "AI 응답은 서버 검증을 통과해야 한다",
     "모델 출력 단계와 애플리케이션 저장 단계에서 서로 다른 오류를 막습니다.",
     '''Layer 1: JSON Schema
  required fields
  2 <= choices <= 4
  0 <= answer <= 3
  additionalProperties = false

Layer 2: Java validation
  answer < choices.size()
  source in provided excerpts
  prompt / explanation not blank
  retry <= 3''',
     [("형식", "구문과 필수 필드를 모델 출력 단계에서 제한", MINT),
      ("의미 범위", "실제 선택지 길이와 정답 인덱스 비교", BLUE),
      ("출처", "이번 요청에서 제공한 경로인지 확인", YELLOW),
      ("회복", "실패 원인을 폐기하고 제한 횟수만 재생성", PINK)],
     "Structured Outputs Schema와 QuizGameService.validateGeneratedQuestion"),
    ("생성 문제의 데이터 계보를 남기는 방법", "설명과 출처까지 저장해 서비스 자산으로",
     "문제만 저장하지 않고 검증과 재사용에 필요한 메타데이터를 함께 보존합니다.",
     '''QuizQuestion
  id
  category / type
  prompt
  choices[]
  answerIndex
  explanation
  source
  signature (unique)
  createdAt

generation -> validation -> save -> reuse''',
     [("재현", "어떤 source에서 생성됐는지 저장", MINT),
      ("설명", "정답과 함께 사용자 피드백 제공", BLUE),
      ("중복", "정규화 signature의 unique 제약", YELLOW),
      ("연속성", "생성 비활성화 시 저장 문제를 다시 출제", PINK)],
     "QuizQuestionEntity와 signature repository 조회, AI OFF fallback 경로"),
]

detail_insertions = []
for i, (title, target_title, premise, code_text, points, evidence) in enumerate(detail_specs):
    detail = engineering_detail_slide(title, premise, code_text, points, evidence, 60 + i)
    detail_insertions.append((prs.slides._sldIdLst[-1], target_title))

# Communication model: distinguish the implemented transport from the scale-up design.
communication_slide = base("통신 구조: 현재 구현과 확장 경계", "DEEP DIVE · COMMUNICATION", 17)
rect(communication_slide, .7, 1.5, 5.65, 4.9, PANEL, LINE)
text(communication_slide, "현재 구현 · HTTP Polling", 1.0, 1.82, 4.9, .35, 19, MINT, True)
rich_lines(communication_slide, [
    ("Transport  ", "TCP/IP 위의 HTTP REST", MINT),
    ("Sync  ", "클라이언트가 1초마다 상태 조회", BLUE),
    ("Command  ", "사용자 행동은 POST 요청으로 즉시 전달", YELLOW),
    ("장점  ", "구현·복구가 단순하고 최대 6명 요구에 충분", MINT),
    ("비용  ", "변화가 없어도 반복 요청과 응답 발생", PINK),
], 1.0, 2.42, 4.9, 3.25, 13)
rect(communication_slide, 6.75, 1.5, 5.85, 4.9, PANEL, LINE)
text(communication_slide, "확장 설계 · WebSocket", 7.05, 1.82, 4.9, .35, 19, BLUE, True)
rich_lines(communication_slide, [
    ("Connection  ", "초기 Upgrade 후 양방향 연결 유지", BLUE),
    ("Push  ", "상태 변경 이벤트를 Room 구독자에게 전파", MINT),
    ("Partition  ", "게임 Room별 채널과 명령 처리 분리", YELLOW),
    ("Scale-out  ", "Redis Pub/Sub로 서버 간 이벤트 공유", PINK),
    ("운영 과제  ", "재연결·순서·중복·backpressure 처리", BLUE),
], 7.05, 2.42, 4.95, 3.25, 13)
text(communication_slide, "포트폴리오 표기 원칙 · 소켓 통신은 확장 설계이며, 현재 동작은 REST 폴링 기반", 1.25, 6.68, 10.8, .28, 14, WHITE, True, align=PP_ALIGN.CENTER)
communication_id = prs.slides._sldIdLst[-1]

# Business and personal closing — placed after the engineering proof.
business_slide = base("비즈니스 모델: 구현 기능에서 출발한 수익 가설", "05 · BUSINESS MODEL", 64)
models = [
    ("BASE", "무료 로컬 버전", "친구·교실이 직접 실행하는 기본 게임과 AI 저장 문제", "목표: 사용 장벽과 배포 마찰 검증", MINT),
    ("SPONSORED", "캠페인 노출", "배너·세션 모달·테마 콘텐츠를 행사 단위로 구성", "목표: 노출보다 참여·클릭·완료 측정", BLUE),
    ("PACKAGE", "수업·행사 패키지", "조직 전용 퀴즈·브랜딩·관리자 운영 화면 제공", "목표: 반복 운영 시간 절감", YELLOW),
    ("HOSTED", "운영형 서비스", "원격 접속 수요가 검증된 뒤 Room·계정·통계를 포함한 호스팅", "목표: 인프라 비용을 감당할 수요 확인", PINK),
]
for i, (tag, head, body, goal, c) in enumerate(models):
    x = .68 + (i % 2) * 6.15
    y = 1.48 + (i // 2) * 2.5
    rect(business_slide, x, y, 5.78, 2.15, PANEL, LINE)
    pill(business_slide, tag, x+.26, y+.26, 1.22, c)
    text(business_slide, head, x+1.72, y+.28, 3.72, .35, 16, WHITE, True)
    text(business_slide, body, x+.26, y+.93, 5.22, .48, 12.2, MUTED)
    text(business_slide, goal, x+.26, y+1.62, 5.22, .28, 11, c, True)
text(business_slide, "현재 상태", .75, 6.55, 1.0, .24, 11, BLUE, True)
text(business_slide, "포트폴리오용 비상업 프로젝트 · 위 항목은 실제 매출이 아닌 검증 순서를 가진 확장 가설", 1.85, 6.48, 10.25, .36, 13, WHITE, True)
business_id = prs.slides._sldIdLst[-1]

business_guard_slide = base("비즈니스 가설을 검증하는 기준", "05 · VALIDATION & GUARDRAILS", 65)
checks = [
    ("VALUE", "사용자가 다시 여는가", "세션당 재경기·7일 내 재사용·진행자 개입 횟수", MINT),
    ("AD", "광고가 경험을 해치지 않는가", "세션당 모달 1회·닫기 가능·배너 레이아웃 안정성", BLUE),
    ("COST", "AI와 운영 비용이 감당되는가", "생성 문제 재사용·fallback 비율·문제당 생성 비용", YELLOW),
    ("RIGHTS", "콘텐츠를 안전하게 쓰는가", "허가된 이미지·원문 링크·대체 텍스트·자료 출처", PINK),
]
for i, (tag, question, metric, c) in enumerate(checks):
    y = 1.5 + i * 1.18
    pill(business_guard_slide, tag, .78, y+.08, 1.0, c)
    text(business_guard_slide, question, 2.08, y+.08, 3.05, .36, 16, WHITE, True)
    rect(business_guard_slide, 5.35, y-.02, 6.95, .82, PANEL, LINE)
    text(business_guard_slide, metric, 5.68, y+.2, 6.35, .36, 13, MUTED)
text(business_guard_slide, "의사결정 원칙", .78, 6.42, 1.35, .24, 11, MINT, True)
text(business_guard_slide, "지표가 개선되지 않으면 기능을 늘리지 않고, 가장 작은 파일럿으로 돌아가 문제와 타깃을 다시 검증합니다.", 2.22, 6.34, 9.92, .45, 14, WHITE, True)
business_guard_id = prs.slides._sldIdLst[-1]

contribution_slide = base("개인 기여: 설계 결정부터 검증까지", "06 · ROLE & CONTRIBUTION", 66)
contributions = [
    ("기획", "LAN·최대 6명·코드 입장이라는 범위를 정의하고 사용자 흐름과 게임 규칙 설계", MINT),
    ("프론트엔드", "React/TypeScript로 게임·채팅·관리자·광고 화면과 애니메이션 구현", BLUE),
    ("백엔드", "Spring Boot로 서버 판정, 상태 머신, 명령 큐, 세션 정리와 API 구현", YELLOW),
    ("데이터·AI", "JPA 트랜잭션과 AI 문제 생성·Schema·출처 검증·DB fallback 구성", PINK),
    ("테스트", "도메인 규칙 JUnit 및 Python 병렬 요청 스크립트와 결과 보고서 작성", MINT),
]
for i, (head, body, c) in enumerate(contributions):
    y = 1.45 + i * 1.02
    text(contribution_slide, f"0{i+1}", .78, y+.2, .52, .25, 12, c, True, "Arial")
    text(contribution_slide, head, 1.52, y+.14, 1.45, .34, 16, WHITE, True)
    rect(contribution_slide, 3.12, y, 9.18, .78, PANEL, LINE)
    text(contribution_slide, body, 3.42, y+.18, 8.6, .38, 13, MUTED)
text(contribution_slide, "기여도 100%", .78, 6.6, 1.4, .25, 12, BLUE, True)
text(contribution_slide, "혼자 구현했기 때문에 ‘무엇을 만들었는가’뿐 아니라 ‘왜 이 경계를 선택했는가’를 코드와 측정 결과로 설명합니다.", 2.3, 6.5, 9.85, .42, 13.5, WHITE, True)
contribution_id = prs.slides._sldIdLst[-1]

future_slide = base("회고에서 다음 검증으로", "06 · RETROSPECTIVE & FUTURE", 67)
future_items = [
    ("배운 점", "공유 상태는 UI보다 서버의 명령 순서·검증·실패 경계에서 신뢰가 결정됩니다.", MINT),
    ("현재 한계", "단일 서버 메모리와 폴링은 현재 6명 범위에는 단순하지만 재시작·수평 확장에 취약합니다.", BLUE),
    ("다음 구현", "Room별 명령 큐, WebSocket push, Redis 상태 외부화와 Controller 통합 테스트를 단계 도입합니다.", YELLOW),
    ("다음 검증", "실제 4~6명 파일럿에서 첫 게임 시작 시간·재경기·오류·진행자 개입을 측정합니다.", PINK),
]
for i, (head, body, c) in enumerate(future_items):
    x = .7 + (i % 2) * 6.17
    y = 1.5 + (i // 2) * 2.35
    rect(future_slide, x, y, 5.8, 1.95, PANEL, LINE)
    text(future_slide, f"0{i+1}", x+.28, y+.28, .48, .25, 12, c, True, "Arial")
    text(future_slide, head, x+.98, y+.23, 1.38, .34, 16, WHITE, True)
    text(future_slide, body, x+.98, y+.8, 4.35, .72, 12.5, MUTED)
text(future_slide, "다음 단계는 더 많은 기능이 아니라, 실제 사용 근거와 실패 복구 범위를 넓히는 것입니다.", .8, 6.5, 11.45, .38, 15, WHITE, True, align=PP_ALIGN.CENTER)
future_id = prs.slides._sldIdLst[-1]

# Keep Result as the final slide after the deep-dive section.
slide_ids = prs.slides._sldIdLst
def slide_id_by_exact_text(value):
    for index, slide in enumerate(prs.slides):
        if any(hasattr(shape, "text") and shape.text.strip() == value for shape in slide.shapes):
            return slide_ids[index]
    raise ValueError(f"slide title not found: {value}")

resolved_qa_insertions = [
    (qa_id, slide_id_by_exact_text(target_title))
    for qa_id, target_title in qa_insertions
]
resolved_detail_insertions = [
    (detail_id, slide_id_by_exact_text(target_title))
    for detail_id, target_title in detail_insertions
]
# Remove movable slides first, then insert them at their narrative positions.
for movable in (
    planning_id, economy_id, concurrency_id, ad_id, load_test_id,
    test_code_id, communication_id, retrospective_id,
    business_id, business_guard_id, contribution_id, future_id, result_id,
    *(qa_id for qa_id, _ in resolved_qa_insertions),
    *(detail_id for detail_id, _ in resolved_detail_insertions),
):
    slide_ids.remove(movable)

def insert_after(movable_id, target_id):
    target_index = list(slide_ids).index(target_id)
    slide_ids.insert(target_index + 1, movable_id)

insert_after(planning_id, slide_id_by_exact_text("프로젝트 소개"))
insert_after(economy_id, slide_id_by_exact_text("System Architecture"))
insert_after(ad_id, slide_id_by_exact_text("공용 채팅과 크레딧 운영 시스템"))
insert_after(communication_id, slide_id_by_exact_text("반복 요청이 만드는 서버 부하"))
insert_after(concurrency_id, slide_id_by_exact_text("명령 큐로 순서와 과부하를 통제"))
insert_after(load_test_id, concurrency_id)
insert_after(test_code_id, load_test_id)

# Place each interview follow-up immediately after its implementation page.
for qa_id, target_id in resolved_qa_insertions:
    target_index = list(slide_ids).index(target_id)
    slide_ids.insert(target_index + 1, qa_id)

# Detail pages are inserted after the evidence page and before its Q&A page.
# Preserve specification order when more than one detail page shares a target.
detail_tail = {}
for detail_id, target_id in resolved_detail_insertions:
    insert_after = detail_tail.get(target_id, target_id)
    target_index = list(slide_ids).index(insert_after)
    slide_ids.insert(target_index + 1, detail_id)
    detail_tail[target_id] = detail_id

# Close the portfolio only after every technical claim, detail and interview
# answer has been shown. This keeps the business section grounded in evidence.
slide_ids.append(business_id)
slide_ids.append(business_guard_id)
slide_ids.append(contribution_id)
slide_ids.append(retrospective_id)
slide_ids.append(future_id)
slide_ids.append(result_id)

# Renumber footers after section reordering so the deck feels manually edited
# and no stale generator index remains.
for page, slide in enumerate(prs.slides, 1):
    for shape in slide.shapes:
        if not hasattr(shape, "text_frame"):
            continue
        if shape.text.startswith("Local Arcade  |"):
            shape.text_frame.paragraphs[0].text = f"Local Arcade  |  {page:02d}"
            p = shape.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.RIGHT
            p.font.name = "Arial"; p.font.size = Pt(8); p.font.color.rgb = MUTED

prs.save(OUT)
print(OUT)
